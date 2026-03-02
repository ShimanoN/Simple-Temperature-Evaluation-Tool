"""
Simple Temperature Evaluation Tool
概要・操作マニュアル PPTX 生成スクリプト
Usage: python create_presentation.py
Output: Simple_Temp_Tool_Overview.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from pptx.util import Emu
import copy
from lxml import etree

# ─── カラーパレット ────────────────────────────────────────
C_NAVY       = RGBColor(0x0D, 0x2B, 0x45)   # ダークネイビー（見出し背景）
C_NAVY_MID   = RGBColor(0x16, 0x45, 0x72)   # ミッドネイビー
C_ORANGE     = RGBColor(0xE8, 0x62, 0x1A)   # M5Stackオレンジ（アクセント）
C_ORANGE_LT  = RGBColor(0xFF, 0x8C, 0x42)   # ライトオレンジ
C_BLUE_LT    = RGBColor(0x00, 0xA8, 0xE8)   # アクアブルー
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BG   = RGBColor(0xF0, 0xF4, 0xF8)   # スライド背景（薄水色）
C_GRAY_LT    = RGBColor(0xE8, 0xEE, 0xF4)   # テーブル偶数行・ボックス
C_GRAY_MID   = RGBColor(0x88, 0x99, 0xAA)   # 補助テキスト
C_DARK_TXT   = RGBColor(0x1A, 0x1A, 0x2E)   # 本文テキスト
C_GREEN      = RGBColor(0x27, 0xAE, 0x60)   # 合格・OK
C_RED        = RGBColor(0xC0, 0x39, 0x2B)   # 警告・バグ

FONT_JP     = "游ゴシック"      # Windows 10/11 標準日本語ゴシック
FONT_JP_B   = "游ゴシック"      # 太字は bold=True で対応
FONT_EN     = "Calibri"

SW  = Inches(13.33)   # スライド幅（16:9）
SH  = Inches(7.5)     # スライド高さ

# ─── ユーティリティ関数 ─────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs


def blank_slide(prs):
    """完全に空白のスライドを追加"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def rgb_fill(shape, color: RGBColor):
    """シェイプの塗りつぶし色を設定"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_border(shape):
    """シェイプの線を非表示"""
    shape.line.fill.background()


def add_rect(slide, left, top, w, h, fill_color=None, border_color=None, border_width_pt=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, w, h
    )
    if fill_color:
        rgb_fill(shape, fill_color)
    else:
        shape.fill.background()
    if border_color and border_width_pt > 0:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width_pt)
    else:
        no_border(shape)
    return shape


def add_textbox(slide, text, left, top, w, h,
                font_name=FONT_JP, font_size=18, bold=False, italic=False,
                color=C_DARK_TXT, align=PP_ALIGN.LEFT,
                v_anchor=None, wrap=True):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_name=FONT_JP, font_size=16, bold=False,
             color=C_DARK_TXT, align=PP_ALIGN.LEFT, space_before_pt=0):
    from pptx.util import Pt as PT2
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = PT2(space_before_pt)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = PT2(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_background(slide, color: RGBColor = C_LIGHT_BG):
    """スライド背景色を設定"""
    add_rect(slide, 0, 0, SW, SH, fill_color=color)


def header_bar(slide, title_text, subtitle_text=""):
    """上部ヘッダーバー（ネイビー）"""
    BAR_H = Inches(1.15)
    # ネイビー帯
    bar = add_rect(slide, 0, 0, SW, BAR_H, fill_color=C_NAVY)
    # 左オレンジアクセントライン
    add_rect(slide, 0, 0, Inches(0.08), BAR_H, fill_color=C_ORANGE)
    # タイトルテキスト
    add_textbox(slide, title_text,
                Inches(0.22), Inches(0.15), Inches(11), Inches(0.65),
                font_size=28, bold=True, color=C_WHITE)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.22), Inches(0.75), Inches(11), Inches(0.35),
                    font_size=14, color=C_ORANGE_LT)


def footer_bar(slide, page_num: int, total: int = 17):
    """下部フッターバー"""
    FH = Inches(0.35)
    add_rect(slide, 0, SH - FH, SW, FH, fill_color=C_NAVY)
    # プロジェクト名（左）
    add_textbox(slide, "Simple Temperature Evaluation Tool  v1.0.0",
                Inches(0.2), SH - FH, Inches(9), FH,
                font_size=10, color=C_GRAY_LT)
    # ページ番号（右）
    add_textbox(slide, f"{page_num} / {total}",
                Inches(11.8), SH - FH, Inches(1.3), FH,
                font_size=10, color=C_GRAY_LT, align=PP_ALIGN.RIGHT)


def section_divider(prs, section_no: str, section_title: str, page_num: int):
    """セクション区切りスライド"""
    slide = blank_slide(prs)
    # フル背景（ネイビー）
    add_rect(slide, 0, 0, SW, SH, fill_color=C_NAVY)
    # 左オレンジ縦ライン
    add_rect(slide, Inches(1.2), Inches(2.2), Inches(0.1), Inches(3.1), fill_color=C_ORANGE)
    # セクション番号
    add_textbox(slide, section_no,
                Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                font_size=72, bold=True, color=C_ORANGE, align=PP_ALIGN.LEFT)
    # セクションタイトル
    add_textbox(slide, section_title,
                Inches(1.5), Inches(3.1), Inches(10), Inches(1.2),
                font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    # フッター
    FH = Inches(0.35)
    add_rect(slide, 0, SH - FH, SW, FH, fill_color=RGBColor(0x05, 0x18, 0x28))
    add_textbox(slide, f"{page_num} / 17",
                Inches(11.8), SH - FH, Inches(1.3), FH,
                font_size=10, color=C_GRAY_LT, align=PP_ALIGN.RIGHT)
    return slide


def info_card(slide, left, top, w, h, label, value,
              label_size=12, value_size=22,
              bg_color=C_WHITE, accent_color=C_NAVY):
    """KPI カードシェイプ（数値強調ボックス）"""
    # 外枠
    card = add_rect(slide, left, top, w, h, fill_color=bg_color)
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    # 上部アクセントライン
    add_rect(slide, left, top, w, Inches(0.07), fill_color=accent_color)
    # ラベル
    add_textbox(slide, label,
                left + Inches(0.15), top + Inches(0.12),
                w - Inches(0.3), Inches(0.35),
                font_size=label_size, color=C_GRAY_MID, align=PP_ALIGN.LEFT)
    # 値
    add_textbox(slide, value,
                left + Inches(0.1), top + Inches(0.4),
                w - Inches(0.2), h - Inches(0.5),
                font_size=value_size, bold=True, color=accent_color,
                align=PP_ALIGN.LEFT)


def make_table(slide, left, top, w, rows_data,
               col_widths=None, header_bg=C_NAVY,
               row_height=Inches(0.42), font_size=14):
    """
    rows_data: list of list of str  （1行目がヘッダー）
    col_widths: list of Inches
    """
    from pptx.util import Pt as PT3
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    total_h = row_height * n_rows

    if col_widths is None:
        unit = w / n_cols
        col_widths = [unit] * n_cols

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top,
                                  w, total_h).table
    # 列幅
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            # 背景
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
                txt_color = C_WHITE
                is_bold = True
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_GRAY_LT
                txt_color = C_DARK_TXT
                is_bold = False
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
                txt_color = C_DARK_TXT
                is_bold = False

            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_JP
                    run.font.size = PT3(font_size)
                    run.font.bold = is_bold
                    run.font.color.rgb = txt_color

    return tbl


# ─── スライド生成関数 ────────────────────────────────────────

def slide_01_title(prs):
    """スライド 1: タイトル"""
    slide = blank_slide(prs)
    # 上下2トーン背景
    add_rect(slide, 0, 0, SW, Inches(4.8), fill_color=C_NAVY)
    add_rect(slide, 0, Inches(4.8), SW, Inches(2.7), fill_color=C_LIGHT_BG)
    # 左垂直アクセント
    add_rect(slide, 0, 0, Inches(0.12), Inches(4.8), fill_color=C_ORANGE)
    # バージョンバッジ
    badge = add_rect(slide, Inches(0.35), Inches(0.25),
                     Inches(1.4), Inches(0.42),
                     fill_color=C_ORANGE)
    badge.line.fill.background()
    add_textbox(slide, "v1.0.0",
                Inches(0.35), Inches(0.25), Inches(1.4), Inches(0.42),
                font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # メインタイトル
    add_textbox(slide, "Simple Temperature",
                Inches(0.35), Inches(0.85), Inches(12.6), Inches(1.0),
                font_size=52, bold=True, color=C_WHITE)
    add_textbox(slide, "Evaluation Tool",
                Inches(0.35), Inches(1.75), Inches(12.6), Inches(1.0),
                font_size=52, bold=True, color=C_ORANGE)
    # サブタイトル
    add_textbox(slide, "概要 ・ 操作マニュアル",
                Inches(0.35), Inches(2.85), Inches(12.6), Inches(0.6),
                font_size=22, color=C_GRAY_LT)
    # 区切り線
    add_rect(slide, Inches(0.35), Inches(3.58), Inches(12.6), Inches(0.03),
             fill_color=C_ORANGE)
    # 概要テキスト（白背景部）
    add_textbox(slide, "ESP32（M5Stack）を使用した温度計測・評価ツール",
                Inches(0.5), Inches(5.0), Inches(10), Inches(0.5),
                font_size=18, bold=True, color=C_NAVY)
    add_textbox(slide,
                "K型熱電対による高精度計測・MicroSD自動記録・リアルタイム統計演算を実現",
                Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.45),
                font_size=14, color=C_DARK_TXT)
    # 日付
    add_textbox(slide, "2026年3月",
                Inches(0.5), Inches(6.15), Inches(4), Inches(0.4),
                font_size=13, color=C_GRAY_MID)
    # フッター
    FH = Inches(0.35)
    add_rect(slide, 0, SH - FH, SW, FH, fill_color=RGBColor(0x05, 0x18, 0x28))
    add_textbox(slide, "1 / 17",
                Inches(11.8), SH - FH, Inches(1.3), FH,
                font_size=10, color=C_GRAY_LT, align=PP_ALIGN.RIGHT)
    return slide


def slide_02_toc(prs):
    """スライド 2: 目次"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "目次", "本資料の構成")
    footer_bar(slide, 2)

    sections = [
        ("Section 1", "プロジェクト概要",     "背景・課題 → ハードウェア → 機能一覧",       "3〜5"),
        ("Section 2", "操作マニュアル",       "起動 → 計測 → 結果確認 → CSV活用",          "6〜9"),
        ("Section 3", "技術詳細",            "アーキテクチャ → タスク設計 → アルゴリズム", "10〜12"),
        ("Section 4", "品質・実績・将来計画", "デバッグ → テスト → コスト → 拡張計画",      "13〜17"),
    ]

    for i, (sec_no, sec_title, sec_desc, pages) in enumerate(sections):
        row = i // 2
        col = i % 2
        cx = Inches(0.4 + col * 6.45)
        cy = Inches(1.35 + row * 2.55)
        cw = Inches(6.15)
        ch = Inches(2.3)
        # カード背景
        card = add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        card.line.color.rgb = C_NAVY
        card.line.width = Pt(1)
        # 左帯
        add_rect(slide, cx, cy, Inches(0.1), ch, fill_color=C_ORANGE)
        # セクション番号
        add_textbox(slide, sec_no,
                    cx + Inches(0.2), cy + Inches(0.1),
                    cw - Inches(0.3), Inches(0.38),
                    font_size=12, color=C_ORANGE, bold=True)
        # セクションタイトル
        add_textbox(slide, sec_title,
                    cx + Inches(0.2), cy + Inches(0.42),
                    cw - Inches(0.3), Inches(0.55),
                    font_size=20, bold=True, color=C_NAVY)
        # 説明
        add_textbox(slide, sec_desc,
                    cx + Inches(0.2), cy + Inches(0.97),
                    cw - Inches(0.3), Inches(0.65),
                    font_size=13, color=C_DARK_TXT)
        # ページ番号
        add_textbox(slide, f"p.{pages}",
                    cx + Inches(0.2), cy + Inches(1.65),
                    cw - Inches(0.3), Inches(0.38),
                    font_size=12, color=C_GRAY_MID)
    return slide


def slide_03_background(prs):
    """スライド 3: 背景と課題"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "背景と課題", "Section 1  プロジェクト概要")
    footer_bar(slide, 3)

    problems = [
        ("📋", "手書き記録",
         "計測値を逐一手書き\nまたはExcelへ手入力"),
        ("📊", "後処理の工数",
         "平均・標準偏差・最大/最小は\n評価終了後に別途計算が必要"),
        ("🔔", "異常検知の遅れ",
         "上限・下限超過の瞬間に\n気づく手段がなかった"),
        ("📁", "記録の不統一",
         "担当者・回によって\n記録フォーマットがばらつく"),
    ]

    # リード文
    add_textbox(slide,
                "温度評価作業（ヒートガン・恒温槽・チャンバー使用）において、次の4つの非効率が存在していました。",
                Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.5),
                font_size=15, color=C_DARK_TXT)

    for i, (icon, title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        cx = Inches(0.4 + col * 6.45)
        cy = Inches(1.85 + row * 2.35)
        cw, ch = Inches(6.15), Inches(2.1)
        card = add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        card.line.color.rgb = RGBColor(0xCC, 0xD5, 0xDE)
        card.line.width = Pt(1)
        add_rect(slide, cx, cy, cw, Inches(0.07), fill_color=C_ORANGE)
        # アイコン + タイトル
        add_textbox(slide, f"{icon}  {title}",
                    cx + Inches(0.15), cy + Inches(0.12),
                    cw - Inches(0.3), Inches(0.5),
                    font_size=17, bold=True, color=C_NAVY)
        # 説明
        add_textbox(slide, desc,
                    cx + Inches(0.15), cy + Inches(0.65),
                    cw - Inches(0.3), Inches(1.2),
                    font_size=14, color=C_DARK_TXT)

    # まとめ帯
    add_rect(slide, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.6),
             fill_color=C_NAVY)
    add_textbox(slide,
                "「記録ミスのリスク」「評価1件あたりの工数」「データ品質の安定性」に直結する課題  →  専用ツールで解決",
                Inches(0.6), Inches(6.55), Inches(12), Inches(0.6),
                font_size=13, color=C_WHITE)
    return slide


def slide_04_hardware(prs):
    """スライド 4: ハードウェア構成"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "ハードウェア構成", "Section 1  プロジェクト概要")
    footer_bar(slide, 4)

    # KPIカード（右側）
    info_card(slide, Inches(9.8), Inches(1.3), Inches(3.1), Inches(1.0),
              "総コスト", "16,082 円", value_size=24, accent_color=C_ORANGE)
    info_card(slide, Inches(9.8), Inches(2.45), Inches(3.1), Inches(1.0),
              "部品点数", "5 部品", value_size=24, accent_color=C_NAVY)
    info_card(slide, Inches(9.8), Inches(3.6), Inches(3.1), Inches(1.0),
              "組み立て時間", "約 15 分", value_size=24, accent_color=C_BLUE_LT)
    info_card(slide, Inches(9.8), Inches(4.75), Inches(3.1), Inches(1.0),
              "ハンダ付け", "不 要", value_size=24, accent_color=C_GREEN)

    # 部品テーブル
    rows = [
        ["部品名",              "型番 / 規格",         "役割",                 "単価"],
        ["M5Stack Core2",      "ESP32-D0WDQ6-V3",    "マイコン本体・LCD・ボタン", "¥6,050"],
        ["K型熱電対センサ",     "Type-K / −200〜1350℃","温度計測",             "¥3,080"],
        ["MAX31855 アンプ基板", "Adafruit #3263",     "熱電対→SPIデジタル変換", "¥2,970"],
        ["MicroSD カード",     "FAT32 / 16GB",        "CSV自動保存",          "¥1,500"],
        ["USB-C ケーブル",     "給電用",              "電源供給",             "¥2,482"],
    ]
    make_table(slide, Inches(0.35), Inches(1.3), Inches(9.2), rows,
               col_widths=[Inches(2.5), Inches(2.4), Inches(2.6), Inches(1.7)],
               font_size=13)

    # 接続メモ
    add_textbox(slide,
                "接続方式：MAX31855 → SPI（GPIO 5/18/19/23）　EEPROM → I2C（GPIO 21/22）",
                Inches(0.35), Inches(5.65), Inches(9.2), Inches(0.4),
                font_size=12, color=C_GRAY_MID)
    # 合計行風バナー
    add_rect(slide, Inches(0.35), Inches(6.1), Inches(9.2), Inches(0.45),
             fill_color=C_NAVY)
    add_textbox(slide, "合計：16,082円（部品実費のみ・外注費 0円）",
                Inches(0.55), Inches(6.1), Inches(9), Inches(0.45),
                font_size=13, bold=True, color=C_WHITE)
    return slide


def slide_05_features(prs):
    """スライド 5: 主要機能"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "主要機能", "Section 1  プロジェクト概要")
    footer_bar(slide, 5)

    features = [
        ("① リアルタイム温度計測",
         "K型熱電対 + MAX31855 使用\n計測範囲：−200〜1350℃  精度：±1℃\n500ms間隔でセンサ読取、LPF（α=0.1）でノイズ除去"),
        ("② 統計値の自動計算・表示",
         "計測中に 平均・標準偏差・最大・最小 を自動更新\nWelford法（逐次更新）でRAMを消費せず処理\n計測終了時に最終統計値をページ表示"),
        ("③ CSVファイル自動保存",
         "MicroSDカードに DATA_0000.csv 形式で自動記録\n10列フォーマット（経過時間・温度・統計・アラームフラグ等）\n計測ごとに連番で自動採番"),
        ("④ アラーム機能",
         "上限（HI）・下限（LO）温度を EEPROM に設定・保存\n超過の瞬間にビープ音 + 画面でリアルタイム通知\nヒステリシス（±5℃）でチラつき防止"),
        ("⑤ 3ボタン操作",
         "BtnA：計測開始/終了/画面遷移\nBtnB：アラーム設定モード進入/設定値+5℃\nBtnC：設定値 −5℃"),
    ]

    for i, (title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        cx = Inches(0.35 + col * 4.33)
        if row == 0:
            cy = Inches(1.3)
            ch = Inches(2.35)
        else:
            cx = Inches(0.35 + (i - 3) * 6.45)
            cy = Inches(3.85)
            ch = Inches(2.35)
            if i >= 5:
                break

        cw = Inches(4.1) if row == 0 else Inches(6.2)

        card = add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        card.line.color.rgb = C_NAVY
        card.line.width = Pt(1)
        add_rect(slide, cx, cy, cw, Inches(0.07), fill_color=C_ORANGE)
        add_textbox(slide, title,
                    cx + Inches(0.12), cy + Inches(0.12),
                    cw - Inches(0.24), Inches(0.5),
                    font_size=14, bold=True, color=C_NAVY)
        add_textbox(slide, desc,
                    cx + Inches(0.12), cy + Inches(0.68),
                    cw - Inches(0.24), ch - Inches(0.8),
                    font_size=12, color=C_DARK_TXT)
    return slide


def slide_06_manual_step1(prs):
    """スライド 6: 操作マニュアル – 起動とIDLE状態"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "操作マニュアル ① — 起動と待機（IDLE）", "Section 2  操作マニュアル")
    footer_bar(slide, 6)

    add_textbox(slide, "Step 1  :  電源を入れる",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.55),
                font_size=22, bold=True, color=C_NAVY)

    rows = [
        ["確認項目",       "表示・動作",                                           "補足"],
        ["起動画面",       '"Temperature Eval Tool  Refactored Version" が数秒表示', "自動的にIDLE画面へ遷移"],
        ["温度表示",       "現在の温度（℃）がリアルタイム表示（500ms更新）",        "センサ未接続時は \"---.- C\" 表示"],
        ["ボタンガイド",   "画面下部に [BtnA] Start  [BtnB] Setting を表示",        "常時表示"],
        ["SDカード状態",   "SDReady または SD Error が画面に表示される",             "初期化失敗時はシリアルにも出力"],
        ["センサ初期化",   "起動時に最大5回リトライ（間隔500ms）",                  "失敗しても起動継続（---.- C表示）"],
        ["センサ異常",     "\"SENSOR ERROR\" が表示される場合",                    "接続・コネクタを確認して再起動"],
    ]
    make_table(slide, Inches(0.35), Inches(1.9), Inches(12.6), rows,
               col_widths=[Inches(2.2), Inches(6.2), Inches(4.2)],
               font_size=13, row_height=Inches(0.5))

    add_textbox(slide,
                "💡  センサが正常接続されていれば起動後すぐに温度表示が始まります。BtnA を押すまで計測データは記録されません。",
                Inches(0.35), Inches(6.45), Inches(12.5), Inches(0.65),
                font_size=13, color=C_NAVY)
    return slide


def slide_07_manual_step2(prs):
    """スライド 7: 操作マニュアル – 計測開始/終了（4状態図付き）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "操作マニュアル ② — 状態遷移と計測操作", "Section 2  操作マニュアル")
    footer_bar(slide, 7)

    # 4状態カード
    states = [
        ("IDLE",          "待機中\n温度リアルタイム表示",
         "BtnA → RUN\nBtnB → ALARM_SETTING",  C_NAVY),
        ("RUN",           "計測中\nCSV自動記録・統計更新",
         "BtnA → RESULT",                      C_ORANGE),
        ("RESULT",        "結果表示\n統計値を2ページ表示",
         "BtnA → IDLE\nBtnB → ページ切替",     C_GREEN),
        ("ALARM\nSETTING", "アラーム設定中\nHI/LO値を調整",
         "BtnA → 項目切替 / 保存→IDLE\nBtnB → +5℃  BtnC → −5℃",
         RGBColor(0x7B, 0x27, 0x82)),
    ]

    for i, (state, desc, btns, col) in enumerate(states):
        cx = Inches(0.35 + i * 3.25)
        cy = Inches(1.3)
        cw, ch = Inches(3.0), Inches(3.2)
        card = add_rect(slide, cx, cy, cw, ch, fill_color=col)
        no_border(card)
        add_textbox(slide, state,
                    cx + Inches(0.12), cy + Inches(0.1),
                    cw - Inches(0.24), Inches(0.85),
                    font_size=19, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    cx + Inches(0.12), cy + Inches(0.95),
                    cw - Inches(0.24), Inches(0.8),
                    font_size=12, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, cx + Inches(0.12), cy + Inches(1.78),
                 cw - Inches(0.24), Inches(0.02),
                 fill_color=RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, btns,
                    cx + Inches(0.12), cy + Inches(1.85),
                    cw - Inches(0.24), Inches(1.25),
                    font_size=11, color=C_WHITE, align=PP_ALIGN.LEFT)
        if i < 3:
            add_textbox(slide, "▶",
                        cx + cw + Inches(0.08), cy + Inches(1.2),
                        Inches(0.15), Inches(0.8),
                        font_size=20, color=C_GRAY_MID, align=PP_ALIGN.CENTER)

    rows = [
        ["状態",            "BtnA",             "BtnB",                   "BtnC",      "CSV記録",  "統計更新"],
        ["IDLE",           "→ RUN（計測開始）", "→ ALARM_SETTING進入",   "なし",      "なし",     "なし"],
        ["RUN",            "→ RESULT（終了）",  "なし",                  "なし",      "10行毎",   "50ms毎"],
        ["RESULT",         "→ IDLE（リセット）","ページ切替（0↔1）",     "なし",      "完了済",   "固定"],
        ["ALARM_SETTING",  "HI→LO / LO保存→IDLE","+5℃ 調整",           "−5℃ 調整", "なし",     "なし"],
    ]
    make_table(slide, Inches(0.35), Inches(4.7), Inches(12.6), rows,
               col_widths=[Inches(2.1), Inches(2.5), Inches(2.7), Inches(1.6), Inches(1.6), Inches(2.1)],
               font_size=12, row_height=Inches(0.44))

    add_textbox(slide,
                "⚠️  RUN → IDLE への直接遷移はなし。必ず RESULT を経由します。ALARM_SETTING はIDLE起点のみ進入可能。",
                Inches(0.35), Inches(6.48), Inches(12.5), Inches(0.5),
                font_size=13, color=C_RED, bold=True)
    return slide


def slide_08_manual_alarm(prs):
    """スライド 8: アラーム設定方法（最新実装準拠）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "操作マニュアル ③ — アラーム設定（ALARM_SETTING）", "Section 2  操作マニュアル")
    footer_bar(slide, 8)

    add_textbox(slide,
                "IDLE状態から BtnB を押すとアラーム設定モード（ALARM_SETTING）に遷移します。設定値は EEPROM に保存され電源OFF後も保持されます。",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.52),
                font_size=14, color=C_DARK_TXT)

    # 設定手順テーブル
    rows = [
        ["操作ステップ",         "ボタン",  "手順の詳細"],
        ["① 設定モードへ進入",   "BtnB",   "IDLE状態で BtnB を押す → ALARM_SETTING画面に遷移（HI設定から開始）"],
        ["② HI温度を調整",      "BtnB/C", "BtnB : HI設定値を +5℃  /  BtnC : HI設定値を −5℃"],
        ["③ LO設定へ切替",      "BtnA",   "HI設定中に BtnA → LO設定に切替（HI値確定）"],
        ["④ LO温度を調整",      "BtnB/C", "BtnB : LO設定値を +5℃  /  BtnC : LO設定値を −5℃"],
        ["⑤ 保存・終了",        "BtnA",   "LO設定中に BtnA → EEPROM保存 → アラームフラグリセット → IDLE に戻る"],
    ]
    make_table(slide, Inches(0.35), Inches(1.82), Inches(12.6), rows,
               col_widths=[Inches(2.6), Inches(1.3), Inches(8.7)],
               font_size=13, row_height=Inches(0.52))

    # アラーム動作説明
    add_textbox(slide, "アラーム発報時の動作",
                Inches(0.35), Inches(4.85), Inches(12), Inches(0.42),
                font_size=16, bold=True, color=C_NAVY)

    alarm_rows = [
        ["状態",                  "画面表示",                     "スピーカー",                "CSVへの記録"],
        ["温度 ≥ HI設定値",       "温度を赤色（RED）で表示",       "2kHz ビープ音（500ms）",    "HI_ALARM 列が 1（true）"],
        ["温度 ≤ LO設定値",       "温度を赤色（RED）で表示",       "1kHz ビープ音（500ms）",    "LO_ALARM 列が 1（true）"],
        ["閾値に近い（ヒステリシス内）", "赤色のまま（クリアされない）", "発報なし（継続監視）",       "フラグ継続"],
        ["通常範囲",              "白色（WHITE）で表示",           "なし",                      "HI_ALARM / LO_ALARM = 0"],
    ]
    make_table(slide, Inches(0.35), Inches(5.35), Inches(12.6), alarm_rows,
               col_widths=[Inches(2.7), Inches(2.8), Inches(2.9), Inches(4.2)],
               font_size=12, row_height=Inches(0.44))

    add_textbox(slide,
                "💡  デフォルト設定 HI = 600℃ / LO = 400℃  /  ヒステリシス幅 = ±5℃（チラつき防止）",
                Inches(0.35), Inches(7.05), Inches(12.5), Inches(0.42),
                font_size=13, color=C_NAVY)
    return slide


def slide_09_csv(prs):
    """スライド 9: CSV出力とExcel活用（最新列定義準拠）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "操作マニュアル ④ — CSV出力とExcel活用", "Section 2  操作マニュアル")
    footer_bar(slide, 9)

    # CSVフォーマット（実際のヘッダ行と一致）
    add_textbox(slide, "CSV ファイルのフォーマット（全10列）",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.42),
                font_size=17, bold=True, color=C_NAVY)
    add_textbox(slide,
                "ヘッダ行:  ElapsedSec, Temp_C, State, Samples, Average_C, StdDev_C, Max_C, Min_C, HI_ALARM, LO_ALARM",
                Inches(0.35), Inches(1.65), Inches(12.5), Inches(0.35),
                font_size=11, color=C_GRAY_MID, italic=True)

    rows = [
        ["列名",        "内容",                              "例"],
        ["ElapsedSec",  "RUN開始からの経過時間（秒）",          "0, 10, 20 …"],
        ["Temp_C",      "LPFフィルタ後の計測温度（℃）",        "125.3"],
        ["State",       "計測中の状態名",                      "RUN"],
        ["Samples",     "累積サンプル数（Welford のカウント）",  "14"],
        ["Average_C",   "その時点の累積平均（℃）",              "124.8"],
        ["StdDev_C",    "その時点の標準偏差（℃）",              "0.42"],
        ["Max_C",       "計測開始からの最大値（℃）",            "128.0"],
        ["Min_C",       "計測開始からの最小値（℃）",            "121.5"],
        ["HI_ALARM",    "上限アラーム発報フラグ",               "0 / 1"],
        ["LO_ALARM",    "下限アラーム発報フラグ",               "0 / 1"],
    ]
    make_table(slide, Inches(0.35), Inches(2.1), Inches(7.8), rows,
               col_widths=[Inches(1.9), Inches(3.8), Inches(2.1)],
               font_size=12, row_height=Inches(0.44))

    # PC側の作業手順
    add_textbox(slide, "PCでの読み込み手順（Excel）",
                Inches(8.35), Inches(1.2), Inches(4.8), Inches(0.42),
                font_size=16, bold=True, color=C_NAVY)

    steps = [
        "① BtnA で計測終了（RESULT画面）",
        "② MicroSDカードをPCに挿入",
        "③ DATA_0000.csv を開く（連番で採番）",
        "④ Excelが列を自動認識して読込み",
        "⑤ ElapsedSec を横軸にグラフ挿入",
        "⑥ HI_ALARM/LO_ALARM列で\n　フィルタをかけて異常点を確認",
    ]
    for j, s in enumerate(steps):
        add_textbox(slide, s,
                    Inches(8.35), Inches(1.72 + j * 0.72),
                    Inches(4.8), Inches(0.65),
                    font_size=13, color=C_DARK_TXT)

    # 注意事項
    add_rect(slide, Inches(0.35), Inches(6.48), Inches(12.6), Inches(0.55),
             fill_color=C_GRAY_LT)
    add_textbox(slide,
                "📁  ファイル名は DATA_0000.csv から自動採番（起動後に /DATA_0000, /DATA_0001 … と増加）。同名上書きなし。",
                Inches(0.5), Inches(6.52), Inches(12.2), Inches(0.48),
                font_size=13, color=C_DARK_TXT)
    return slide


def slide_10_architecture(prs):
    """スライド 10: ソフトウェアアーキテクチャ（PLC→C++転用）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "ソフトウェアアーキテクチャ", "Section 3  技術詳細")
    footer_bar(slide, 10)

    add_textbox(slide, "PLC 設計経験 → C++ 組み込み開発への体系的転用",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.5),
                font_size=17, bold=True, color=C_NAVY)

    rows = [
        ["PLC の設計手法",        "C++ での実装",                            "目的・効果"],
        ["スキャン周期（固定周期）",  "millis() による3タスク時分割\nIO:10ms / Logic:50ms / UI:200ms", "センサ読取と画面描画の干渉を防止"],
        ["グローバルラベル\n（全局共有変数）", "struct G { … } による\n全変数の一元管理",          "モジュール間の変数参照を一箇所に集約"],
        ["工程歩進（S0→S1→S2）",   "State Machine\nIDLE→RUN→RESULT→IDLE",   "状態に応じた処理切替・誤動作防止"],
        ["インターロック条件",      "HI/LO アラーム判定ロジック",              "設定温度超過のリアルタイム検知"],
    ]
    make_table(slide, Inches(0.35), Inches(1.8), Inches(12.6), rows,
               col_widths=[Inches(3.1), Inches(4.7), Inches(4.8)],
               font_size=13, row_height=Inches(0.62))

    add_textbox(slide,
                "💡  特定の言語・ツールへの依存ではなく「制御の本質」として理解することで、異なる技術環境でも同様の設計が適用可能",
                Inches(0.35), Inches(6.38), Inches(12.5), Inches(0.65),
                font_size=13, color=C_NAVY)
    return slide


def slide_11_tasks(prs):
    """スライド 11: タスク設計（各タスクの実際の処理内容を正確に記載）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "タスク設計（3タスク・マルチレート）", "Section 3  技術詳細")
    footer_bar(slide, 11)

    add_textbox(slide, "単一 loop() を使わず、周期の異なる3タスクを millis() で時分割実行することで処理干渉を防止",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.5),
                font_size=15, color=C_DARK_TXT)

    tasks = [
        ("IO_Task",    "10 ms",  "95%",
         "・M5.update()（ボタンエッジ検出）\n・センサ読取（MAX31855 / SPI）\n　　読取間隔は TC_READ_INTERVAL=500ms\n・LPFフィルタ適用（α=0.1）\n・アラーム判定（ヒステリシス付き）\n・SD書込（10行バッファごと）"),
        ("Logic_Task", "50 ms",  "98%",
         "・Welford法による逐次統計更新\n　（Average / StdDev / Max / Min）\n・ボタンA/B/C イベント処理\n・状態遷移制御\n・EEPROM保存処理（設定確定時）"),
        ("UI_Task",    "200 ms", "97.5%",
         "・状態別LCD画面描画\n　IDLE / RUN / RESULT(2ページ)\n　ALARM_SETTING\n・センサエラー / SD状態表示\n・アラーム時は温度を赤色表示"),
    ]

    for i, (name, period, margin, detail) in enumerate(tasks):
        cx = Inches(0.35 + i * 4.33)
        cy = Inches(1.85)
        cw, ch = Inches(4.1), Inches(4.05)
        card = add_rect(slide, cx, cy, cw, ch, fill_color=C_WHITE)
        card.line.color.rgb = C_NAVY
        card.line.width = Pt(1.5)
        add_rect(slide, cx, cy, cw, Inches(0.08), fill_color=C_ORANGE)
        add_textbox(slide, name,
                    cx + Inches(0.15), cy + Inches(0.12),
                    cw - Inches(0.3), Inches(0.55),
                    font_size=18, bold=True, color=C_NAVY)
        add_textbox(slide, f"周期：{period}",
                    cx + Inches(0.15), cy + Inches(0.68),
                    cw - Inches(0.3), Inches(0.4),
                    font_size=14, color=C_ORANGE, bold=True)
        add_textbox(slide, f"実行余裕率：{margin}",
                    cx + Inches(0.15), cy + Inches(1.05),
                    cw - Inches(0.3), Inches(0.4),
                    font_size=13, color=C_GREEN)
        add_textbox(slide, detail,
                    cx + Inches(0.15), cy + Inches(1.52),
                    cw - Inches(0.3), Inches(2.4),
                    font_size=12, color=C_DARK_TXT)

    add_textbox(slide,
                "💡  IO_Taskのセンサ読取は10ms周期だが実際の読取は TC_READ_INTERVAL_MS=500ms ごと（MAX31855変換完了を待つ設計）",
                Inches(0.35), Inches(6.18), Inches(12.5), Inches(0.38),
                font_size=12, color=C_NAVY)
    return slide


def slide_12_algorithms(prs):
    """スライド 12: アルゴリズム（実際のコードに準拠した式）"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "主要アルゴリズム", "Section 3  技術詳細")
    footer_bar(slide, 12)

    # Welford法
    add_textbox(slide, "① Welford 法（オンライン統計計算）",
                Inches(0.35), Inches(1.2), Inches(12), Inches(0.45),
                font_size=18, bold=True, color=C_NAVY)
    welford = add_rect(slide, Inches(0.35), Inches(1.75), Inches(8.4), Inches(2.45),
                       fill_color=C_WHITE)
    welford.line.color.rgb = C_GRAY_MID
    welford.line.width = Pt(1)
    add_textbox(slide,
                "計測データを配列に蓄積せず、サンプルが来るたびに逐次的に平均・分散・標準偏差を更新。\n"
                "double型で計算することで float の丸め誤差を最小化。\n\n"
                "  delta  = x[n] − prevMean\n"
                "  Sum   += x[n]  →  newMean = Sum / Count\n"
                "  M2    += delta × (x[n] − newMean)\n"
                "  StdDev = √(M2 / Count)   ← Logic_Task 50ms周期で毎回計算",
                Inches(0.5), Inches(1.82), Inches(8.0), Inches(2.3),
                font_size=13, color=C_DARK_TXT)
    info_card(slide, Inches(9.0), Inches(1.75), Inches(3.6), Inches(2.45),
              "RAM 使用率", "7.2%\n（計測時間に\n依存しない）",
              value_size=18, accent_color=C_GREEN)

    # ローパスフィルタ
    add_textbox(slide, "② 一次遅れローパスフィルタ（ノイズ除去）",
                Inches(0.35), Inches(4.35), Inches(12), Inches(0.45),
                font_size=18, bold=True, color=C_NAVY)
    lpf = add_rect(slide, Inches(0.35), Inches(4.9), Inches(8.4), Inches(1.9),
                   fill_color=C_WHITE)
    lpf.line.color.rgb = C_GRAY_MID
    lpf.line.width = Pt(1)
    add_textbox(slide,
                "IO_Task（10ms周期）内で TC_READ_INTERVAL_MS=500ms の間隔で適用。\n"
                "α = FILTER_ALPHA = 0.1f（Global.h 定義）\n\n"
                "  filtered[n] = filtered[n-1] × (1 − α)  +  raw[n] × α\n\n"
                "初回読取時は raw値をそのまま初期値に設定（NaN収束待ち回避）。",
                Inches(0.5), Inches(4.97), Inches(8.0), Inches(1.78),
                font_size=13, color=C_DARK_TXT)
    info_card(slide, Inches(9.0), Inches(4.9), Inches(3.6), Inches(1.9),
              "計測精度", "±1℃\n（MAX31855\n仕様準拠）",
              value_size=18, accent_color=C_ORANGE)

    # ヒステリシスアラーム
    add_textbox(slide, "③ ヒステリシス付きアラーム判定（updateAlarmFlags）",
                Inches(0.35), Inches(6.95), Inches(12), Inches(0.38),
                font_size=15, bold=True, color=C_NAVY)
    return slide


def slide_13_bugs(prs):
    """スライド 13: デバッグ実績"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "デバッグ実績（5件）", "Section 4  品質・実績・将来計画")
    footer_bar(slide, 13)

    add_textbox(slide,
                "ロジックアナライザなし・シリアルモニタ + CSV 目視という環境下で全バグを根本原因まで追跡・修正",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.45),
                font_size=14, color=C_DARK_TXT)

    # 代表事例
    add_textbox(slide, "代表事例：Bug-01 → Bug-02（複合バグ）",
                Inches(0.35), Inches(1.75), Inches(12), Inches(0.45),
                font_size=16, bold=True, color=C_NAVY)

    bug_rows = [
        ["ステップ",         "発見内容・判断"],
        ["① 症状の発見",     "CSV の Average_C 列が全行 0.0 のまま"],
        ["② 仮説立案",       "「書込み条件が早すぎてWelfordが未実行」と仮説"],
        ["③ 根本原因1（Bug-01）", "IO_Task が Logic_Task より先に実行 → 0.0 を書き込んでいた"],
        ["④ 修正1",          "書込み条件を D_Count ≥ 10（100ms待機）に変更"],
        ["⑤ 再テストで再現",  "修正後もまだ 0.0 → 別原因が存在すると判断してコード精読"],
        ["⑥ 根本原因2（Bug-02）", "Average 計算が RESULT 遷移時のみに実装されていた設計バグ"],
        ["⑦ 修正2・完全解決", "Logic_Task（50ms周期）内で毎回計算するよう変更 → 解消"],
    ]
    make_table(slide, Inches(0.35), Inches(2.3), Inches(12.6), bug_rows,
               col_widths=[Inches(3.5), Inches(9.1)],
               font_size=12, row_height=Inches(0.44))

    add_textbox(slide,
                "💡 「症状が同じでも根本原因は別かもしれない」という視点で追跡した複合バグ解析の事例",
                Inches(0.35), Inches(6.5), Inches(12.5), Inches(0.52),
                font_size=13, color=C_NAVY)
    return slide


def slide_14_bugs_list(prs):
    """スライド 14: バグ一覧"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "バグ一覧（全5件）", "Section 4  品質・実績・将来計画")
    footer_bar(slide, 14)

    rows = [
        ["バグ",    "症状",                  "根本原因",                               "解決策"],
        ["Bug-01", "CSV の Average が 0.0",  "IO_Task が Logic_Task より先に実行",      "書込み条件を D_Count ≥ 10 に変更"],
        ["Bug-02", "同上（修正後も再現）",    "Average 計算が RESULT 遷移時のみ",        "Logic_Task 内で毎回計算するよう変更"],
        ["Bug-03", "温度読取の瞬間スパイク",  "SPI 通信の一時的乱れ",                   "最大3回リトライ処理を実装"],
        ["Bug-04", "SD 書込み後の読取値不定", "SPI 共用バスの状態不定",                 "SPI.end(); SPI.begin(); でバスリセット"],
        ["Bug-05", "MicroSD 初期化失敗",      "CS ピン未指定（GPIO4 が明示されていない）", "SD.begin(4) で GPIO4 に明示化"],
    ]
    make_table(slide, Inches(0.35), Inches(1.3), Inches(12.6), rows,
               col_widths=[Inches(1.3), Inches(3.1), Inches(4.5), Inches(3.7)],
               font_size=13, row_height=Inches(0.6))

    # KPIカード
    info_card(slide, Inches(0.4),  Inches(5.65), Inches(4.0), Inches(1.1),
              "解決バグ数", "5 / 5 件（100%）", value_size=22, accent_color=C_GREEN)
    info_card(slide, Inches(4.55), Inches(5.65), Inches(4.0), Inches(1.1),
              "デバッグ手法", "シリアルモニタ + CSV 目視", value_size=14, accent_color=C_NAVY)
    info_card(slide, Inches(8.7),  Inches(5.65), Inches(4.3), Inches(1.1),
              "ドキュメント", "TROUBLESHOOTING.md 722行", value_size=14, accent_color=C_ORANGE)
    return slide


def slide_15_test(prs):
    """スライド 15: テスト実績"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "テスト実績", "Section 4  品質・実績・将来計画")
    footer_bar(slide, 15)

    add_textbox(slide, "7カテゴリ 33ケース の統合テストを全て手動実機テストで実施",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.45),
                font_size=15, color=C_DARK_TXT)

    rows = [
        ["カテゴリ",           "テストケース数", "合格数", "合格率"],
        ["① 起動・初期化",    "4",             "4",      "100%"],
        ["② 温度計測",        "5",             "5",      "100%"],
        ["③ CSV自動保存",     "6",             "6",      "100%"],
        ["④ 統計計算",        "5",             "5",      "100%"],
        ["⑤ アラーム機能",    "5",             "5",      "100%"],
        ["⑥ 状態遷移",        "4",             "4",      "100%"],
        ["⑦ エラー対応",      "4",             "4",      "100%"],
        ["合計",              "33",            "33",     "100%"],
    ]
    make_table(slide, Inches(0.35), Inches(1.8), Inches(7.5), rows,
               col_widths=[Inches(3.0), Inches(1.7), Inches(1.5), Inches(1.3)],
               font_size=14, row_height=Inches(0.5))

    # KPIカード群
    info_card(slide, Inches(8.1), Inches(1.8), Inches(4.8), Inches(1.4),
              "テスト合格率", "33 / 33  ✓  100%", value_size=22, accent_color=C_GREEN)
    info_card(slide, Inches(8.1), Inches(3.4), Inches(4.8), Inches(1.1),
              "Flash 使用率", "31.0%  （余裕 69%）", value_size=18, accent_color=C_NAVY)
    info_card(slide, Inches(8.1), Inches(4.65), Inches(4.8), Inches(1.1),
              "RAM 使用率",   "7.2%  （余裕 92.8%）", value_size=18, accent_color=C_BLUE_LT)
    info_card(slide, Inches(8.1), Inches(5.9), Inches(4.8), Inches(1.1),
              "ドキュメント整備", "15 ファイル（自作）", value_size=18, accent_color=C_ORANGE)
    return slide


def slide_16_expansion(prs):
    """スライド 16: 将来の拡張計画"""
    slide = blank_slide(prs)
    slide_background(slide)
    header_bar(slide, "将来の拡張計画", "Section 4  品質・実績・将来計画")
    footer_bar(slide, 16)

    add_textbox(slide, "Flash余裕率 69% / RAM余裕率 92.8% — 機能追加の余地が充分にあります",
                Inches(0.35), Inches(1.2), Inches(12.5), Inches(0.45),
                font_size=15, color=C_DARK_TXT)

    rows = [
        ["フェーズ", "追加機能",                    "技術的要件・難易度"],
        ["Phase 1\n（短期）",   "タイムスタンプ記録",
         "RTC モジュール（I2C）追加\n難易度：低"],
        ["Phase 2\n（短期）",   "多点同時計測\n（マルチチャンネル）",
         "MAX31855 複数接続（SPI CS 切替）\n難易度：中"],
        ["Phase 3\n（中期）",   "WiFi によるリアルタイム\nデータ送信",
         "ESP32 内蔵WiFi + MQTT / HTTP\n難易度：中"],
        ["Phase 4\n（長期）",   "計測対象の汎用化\n（湿度・圧力・電流）",
         "センサ抽象化レイヤーの設計\n難易度：高"],
    ]
    make_table(slide, Inches(0.35), Inches(1.8), Inches(12.6), rows,
               col_widths=[Inches(1.6), Inches(3.5), Inches(7.5)],
               font_size=13, row_height=Inches(0.95))

    add_textbox(slide,
                "現在のアーキテクチャはタスク分離・モジュール分割が完了しており、機能追加時の影響範囲を最小化できる設計になっています。",
                Inches(0.35), Inches(6.38), Inches(12.5), Inches(0.65),
                font_size=13, color=C_NAVY)
    return slide


def slide_17_summary(prs):
    """スライド 17: まとめ"""
    slide = blank_slide(prs)
    # 背景：ネイビー
    add_rect(slide, 0, 0, SW, SH, fill_color=C_NAVY)
    add_rect(slide, 0, 0, Inches(0.12), SH, fill_color=C_ORANGE)
    # ヘッダーテキスト
    add_textbox(slide, "まとめ",
                Inches(0.35), Inches(0.2), Inches(12), Inches(0.65),
                font_size=30, bold=True, color=C_WHITE)
    add_textbox(slide, "Simple Temperature Evaluation Tool  v1.0.0",
                Inches(0.35), Inches(0.82), Inches(12), Inches(0.45),
                font_size=15, color=C_ORANGE_LT)

    # 区切り線
    add_rect(slide, Inches(0.35), Inches(1.35), Inches(12.6), Inches(0.03),
             fill_color=C_ORANGE)

    points = [
        ("コスト",     "16,082円 / 外注費0円 / 組み立て15分"),
        ("品質",       "33/33テスト合格（100%）/ Flash 31% / RAM 7.2%"),
        ("操作性",     "ワンボタン操作・手順書不要・自動CSV記録"),
        ("設計",       "PLC経験をC++に転用した体系的アーキテクチャ"),
        ("拡張性",     "Flash/RAM余裕率90%超・Phase 4まで拡張ロードマップあり"),
    ]

    for i, (key, val) in enumerate(points):
        cy = Inches(1.55 + i * 0.92)
        add_rect(slide, Inches(0.35), cy, Inches(1.5), Inches(0.7),
                 fill_color=C_ORANGE)
        add_textbox(slide, key,
                    Inches(0.35), cy, Inches(1.5), Inches(0.7),
                    font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        desc_box = add_rect(slide, Inches(1.85), cy, Inches(11.1), Inches(0.7),
                            fill_color=C_NAVY_MID)
        no_border(desc_box)
        add_textbox(slide, val,
                    Inches(2.0), cy, Inches(10.8), Inches(0.7),
                    font_size=15, color=C_WHITE)

    # フッター
    FH = Inches(0.35)
    add_rect(slide, 0, SH - FH, SW, FH, fill_color=RGBColor(0x05, 0x18, 0x28))
    add_textbox(slide, "Simple Temperature Evaluation Tool  v1.0.0  |  2026年3月",
                Inches(0.2), SH - FH, Inches(9), FH,
                font_size=10, color=C_GRAY_LT)
    add_textbox(slide, "17 / 17",
                Inches(11.8), SH - FH, Inches(1.3), FH,
                font_size=10, color=C_GRAY_LT, align=PP_ALIGN.RIGHT)
    return slide


# ─── メイン ──────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("スライド生成中...")

    slide_01_title(prs)
    print("  [1/17] タイトル")

    slide_02_toc(prs)
    print("  [2/17] 目次")

    # Section 1
    slide_03_background(prs)
    print("  [3/17] 背景と課題")

    slide_04_hardware(prs)
    print("  [4/17] ハードウェア構成")

    slide_05_features(prs)
    print("  [5/17] 主要機能")

    # Section 2
    slide_06_manual_step1(prs)
    print("  [6/17] 操作マニュアル① 起動・IDLE")

    slide_07_manual_step2(prs)
    print("  [7/17] 操作マニュアル② RUN/RESULT")

    slide_08_manual_alarm(prs)
    print("  [8/17] 操作マニュアル③ アラーム設定")

    slide_09_csv(prs)
    print("  [9/17] 操作マニュアル④ CSV/Excel")

    # Section 3
    slide_10_architecture(prs)
    print("  [10/17] ソフトウェアアーキテクチャ")

    slide_11_tasks(prs)
    print("  [11/17] タスク設計")

    slide_12_algorithms(prs)
    print("  [12/17] アルゴリズム")

    # Section 4
    slide_13_bugs(prs)
    print("  [13/17] デバッグ代表事例")

    slide_14_bugs_list(prs)
    print("  [14/17] バグ一覧")

    slide_15_test(prs)
    print("  [15/17] テスト実績")

    slide_16_expansion(prs)
    print("  [16/17] 将来の拡張計画")

    slide_17_summary(prs)
    print("  [17/17] まとめ")

    output_path = "Simple_Temp_Tool_Overview.pptx"
    prs.save(output_path)
    print(f"\n✅ 完了: {output_path}")
    print(f"   スライド数: {len(prs.slides)} 枚")


if __name__ == "__main__":
    main()
