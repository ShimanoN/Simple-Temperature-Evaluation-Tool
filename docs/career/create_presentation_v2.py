#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
create_presentation_v2.py
--------------------------
Simple Temperature Evaluation Tool — PowerPoint 自動生成スクリプト v2.0
出力: docs/career/Simple_Temp_Tool_Overview_v2.pptx
仕様書: docs/PPTX_SPEC.md

使い方:
  python create_presentation_v2.py
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# ▼ 色定数
# ---------------------------------------------------------------------------
C_NAVY     = RGBColor(0x0D, 0x47, 0x7F)
C_ORANGE   = RGBColor(0xFF, 0x6D, 0x00)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_GRAY_LT  = RGBColor(0xF4, 0xF4, 0xF4)
C_GRAY_MID = RGBColor(0x99, 0x99, 0x99)
C_PURPLE   = RGBColor(0x7B, 0x27, 0x82)
C_RED_LT   = RGBColor(0xFF, 0xE0, 0xE0)
C_BLUE     = RGBColor(0x1E, 0x90, 0xFF)
C_BLACK    = RGBColor(0x00, 0x00, 0x00)

# ---------------------------------------------------------------------------
# ▼ フォント定数
# ---------------------------------------------------------------------------
FONT_JP   = "游ゴシック"
FONT_EN   = "Segoe UI"
FONT_CODE = "Consolas"

# ---------------------------------------------------------------------------
# ▼ レイアウト定数
# ---------------------------------------------------------------------------
W  = Inches(13.33)
H  = Inches(7.5)
HDR_H  = Inches(1.0)
FTR_Y  = Inches(6.9)
FTR_H  = Inches(0.6)
CX     = Inches(0.35)   # コンテンツエリア X 開始
CY     = Inches(1.1)    # コンテンツエリア Y 開始
CW     = Inches(12.63)  # コンテンツ幅
CH     = Inches(5.75)   # コンテンツ高さ

RECT = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE の整数値

TOTAL_SLIDES = 23

# ---------------------------------------------------------------------------
# ▼ ヘルパー関数
# ---------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    """白背景のブランクスライドを追加して返す"""
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def _set_bg(slide, color: RGBColor):
    """スライド背景を単色塗り"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor,
             line_color=None, line_width=None):
    """塗りつぶし矩形を追加して shape を返す"""
    shape = slide.shapes.add_shape(RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1.0)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                size=Pt(13), color=None, bold=False, italic=False,
                font=FONT_JP, align=PP_ALIGN.LEFT, wrap=True):
    """単一テキストのテキストボックスを追加"""
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_multiline_textbox(slide, lines, x, y, w, h,
                          size=Pt(13), color=None, bold=False, italic=False,
                          font=FONT_JP, align=PP_ALIGN.LEFT, wrap=True):
    """複数行テキストボックスを追加（lines は str のリスト）"""
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
    return tb


def bullet_box(slide, items, x, y, w, h,
               size=Pt(13), color=None, font=FONT_JP, bold=False):
    """箇条書きテキストボックス（各要素に ・ を付ける）"""
    if color is None:
        color = C_DARK_TXT
    text = "\n".join(items)
    return add_multiline_textbox(slide, items, x, y, w, h,
                                 size=size, color=color, bold=bold,
                                 font=font, align=PP_ALIGN.LEFT)


def header_bar(slide, title: str, section: str = ""):
    """ネイビーヘッダーバー + タイトル + セクションラベル"""
    add_rect(slide, Inches(0), Inches(0), W, HDR_H, C_NAVY)
    add_textbox(slide, title,
                Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.55),
                size=Pt(24), color=C_WHITE, bold=True, font=FONT_EN,
                align=PP_ALIGN.LEFT)
    if section:
        add_textbox(slide, section,
                    Inches(0.3), Inches(0.65), Inches(12.5), Inches(0.3),
                    size=Pt(12), color=C_ORANGE, font=FONT_JP,
                    align=PP_ALIGN.LEFT)


def footer_bar(slide, page_num: int):
    """ネイビーフッターバー + プロジェクト名 + ページ番号"""
    add_rect(slide, Inches(0), FTR_Y, W, FTR_H, C_NAVY)
    add_textbox(slide, "Simple Temperature Evaluation Tool",
                Inches(0.3), Inches(6.93), Inches(8.0), Inches(0.35),
                size=Pt(11), color=C_ORANGE, font=FONT_EN,
                align=PP_ALIGN.LEFT)
    add_textbox(slide, str(page_num),
                Inches(11.5), Inches(6.93), Inches(1.5), Inches(0.35),
                size=Pt(11), color=C_WHITE, font=FONT_EN,
                align=PP_ALIGN.RIGHT)


def make_table(slide, x, y, total_w, rows_data, col_widths,
               row_height=Inches(0.46),
               hdr_fill=C_NAVY, hdr_txt=C_WHITE,
               odd_fill=C_WHITE, even_fill=C_GRAY_LT,
               txt_color=C_DARK_TXT,
               font_size=Pt(12), hdr_size=Pt(13),
               code_font=False):
    """
    テーブルを追加。rows_data の 1 行目はヘッダー。
    col_widths: list of Inches (合計値は total_w と一致させること)
    """
    nr = len(rows_data)
    nc = len(col_widths)
    total_h = row_height * nr
    tbl_shape = slide.shapes.add_table(nr, nc, x, y, total_w, total_h)
    tbl = tbl_shape.table

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri in range(nr):
        tbl.rows[ri].height = row_height

    for ri, row in enumerate(rows_data):
        is_hdr = (ri == 0)
        if is_hdr:
            fill = hdr_fill
            fc   = hdr_txt
            fs   = hdr_size
            fn   = FONT_EN
            fb   = True
        else:
            fill = even_fill if ri % 2 == 0 else odd_fill
            fc   = txt_color
            fs   = font_size
            fn   = FONT_CODE if code_font else FONT_JP
            fb   = False

        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(cell_text)
            run.font.size  = fs
            run.font.bold  = fb
            run.font.color.rgb = fc
            run.font.name  = fn

    return tbl_shape


def color_row(tbl_shape, ri, fill_rgb, txt_rgb, bold=True):
    """指定行を上書き着色（make_table 呼び出し後に使用）"""
    tbl = tbl_shape.table
    nc = len(tbl.columns)
    for ci in range(nc):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_rgb
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = txt_rgb
                run.font.bold = bold


def lcd_mock(slide, x, y, w, h, lines, border_color=C_GRAY_MID):
    """黒背景 LCD 風ブロックにモノスペーステキストを表示"""
    add_rect(slide, x, y, w, h, C_BLACK, border_color, line_width=2.0)
    tb = slide.shapes.add_textbox(
        x + Inches(0.12), y + Inches(0.12),
        w - Inches(0.24), h - Inches(0.24)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name  = FONT_CODE
        run.font.size  = Pt(13)
        run.font.color.rgb = C_WHITE


def section_badge(slide, label, x, y):
    """セクションバッジ: 幅0.9"×0.35" ネイビー矩形＋白テキスト"""
    add_rect(slide, x, y, Inches(0.9), Inches(0.35), C_NAVY)
    add_textbox(slide, label,
                x + Inches(0.05), y + Inches(0.02),
                Inches(0.8), Inches(0.32),
                size=Pt(11), color=C_WHITE, bold=True,
                font=FONT_EN, align=PP_ALIGN.CENTER)


def toc_row(slide, label, text, slides_str, x, y):
    """目次の1行: バッジ + テキスト + ページ範囲"""
    section_badge(slide, label, x, y)
    add_textbox(slide, text,
                x + Inches(1.0), y,
                Inches(3.5), Inches(0.35),
                size=Pt(14), color=C_DARK_TXT, font=FONT_JP)
    add_textbox(slide, slides_str,
                x + Inches(4.6), y,
                Inches(1.3), Inches(0.35),
                size=Pt(13), color=C_GRAY_MID, font=FONT_EN,
                align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# ▼ スライド 1 — タイトル
# ---------------------------------------------------------------------------
def slide01(prs):
    s = _blank_slide(prs)
    _set_bg(s, C_NAVY)

    add_textbox(s, "Simple Temperature Evaluation Tool",
                Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.9),
                size=Pt(36), color=C_WHITE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "M5Stack Basic V2.7 × MAX31855  熱電対温度計測システム",
                Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.55),
                size=Pt(18), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)
    # オレンジ区切り線
    add_rect(s, Inches(1.0), Inches(3.62), Inches(11.3), Inches(0.04), C_ORANGE)
    add_textbox(s, "リアルタイム計測 ／ 統計解析 ／ SD カード記録 ／ アラーム機能",
                Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.45),
                size=Pt(15), color=C_ORANGE, font=FONT_JP,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "ESP32  /  C++ (Arduino)  /  PlatformIO",
                Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.45),
                size=Pt(14), color=C_GRAY_MID, font=FONT_EN,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "v1.0.0  ｜  2026 年 3 月",
                Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.4),
                size=Pt(13), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# ▼ スライド 2 — 目次
# ---------------------------------------------------------------------------
def slide02(prs):
    s = _blank_slide(prs)
    header_bar(s, "目次", "Table of Contents")
    footer_bar(s, 2)

    # 左列
    lx = Inches(0.35)
    ly_base = Inches(1.2)
    step = Inches(0.7)
    left_rows = [
        ("S1", "プロジェクト概要",  "Slides  3-5"),
        ("S2", "ハードウェア",      "Slides  6-7"),
        ("S3", "ソフトウェア設計",  "Slides  8-9"),
        ("S4", "操作マニュアル",    "Slides 10-14"),
    ]
    for i, (badge, text, pages) in enumerate(left_rows):
        toc_row(s, badge, text, pages, lx, ly_base + step * i)

    # 右列
    rx = Inches(6.7)
    right_rows = [
        ("S5", "データ管理",    "Slides 15-16"),
        ("S6", "技術詳細",      "Slides 17-20"),
        ("S7", "品質・テスト",  "Slides 21-22"),
        ("Sum", "Summary",       "Slide  23"),
    ]
    for i, (badge, text, pages) in enumerate(right_rows):
        toc_row(s, badge, text, pages, rx, ly_base + step * i)


# ---------------------------------------------------------------------------
# ▼ スライド 3 — 背景・目的
# ---------------------------------------------------------------------------
def slide03(prs):
    s = _blank_slide(prs)
    header_bar(s, "背景と目的", "Section 1  プロジェクト概要")
    footer_bar(s, 3)

    # 背景ブロック
    add_rect(s, Inches(0.35), Inches(1.1), Inches(12.6), Inches(1.6), C_GRAY_LT)
    add_multiline_textbox(s, [
        "・現場での K 型熱電対温度モニタリングを簡単に行いたい",
        "・既製品は機能過多・高価。M5Stack + MAX31855 で必要最小限のツールを自作",
        "・PLC ラダー設計の知識を C++ に落とし込む練習として開発",
    ],
        Inches(0.55), Inches(1.18), Inches(12.2), Inches(1.42),
        size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    # 課題ブロック
    add_rect(s, Inches(0.35), Inches(2.85), Inches(6.0), Inches(2.8),
             C_WHITE, C_RED, line_width=1.5)
    add_textbox(s, "従来の課題",
                Inches(0.5), Inches(2.93), Inches(5.5), Inches(0.4),
                size=Pt(16), color=C_RED, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・手書き記録 → 転記ミス・工数大",
        "・既製ロガーは PC 接続必須で現場には不向き",
        "・データをその場で確認できない",
    ],
        Inches(0.5), Inches(3.4), Inches(5.7), Inches(2.1),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    # 解決ブロック
    add_rect(s, Inches(6.7), Inches(2.85), Inches(6.0), Inches(2.8),
             C_WHITE, C_GREEN, line_width=1.5)
    add_textbox(s, "本ツールで解決",
                Inches(6.85), Inches(2.93), Inches(5.5), Inches(0.4),
                size=Pt(16), color=C_GREEN, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・MicroSD に自動 CSV 記録（手書き不要）",
        "・スタンドアロン動作（PC 不要）",
        "・LCD でリアルタイム確認 + アラーム通知",
    ],
        Inches(6.85), Inches(3.4), Inches(5.7), Inches(2.1),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 4 — コストメリット比較
# ---------------------------------------------------------------------------
def slide04(prs):
    s = _blank_slide(prs)
    header_bar(s, "コストメリット — 自作 vs. 市販品", "Section 1  プロジェクト概要")
    footer_bar(s, 4)

    # 左ブロック見出し
    add_textbox(s, "本ツール 品目別コスト",
                Inches(0.35), Inches(1.1), Inches(5.8), Inches(0.35),
                size=Pt(15), color=C_NAVY, bold=True, font=FONT_JP)

    # コスト内訳テーブル
    cost_rows = [
        ["部品", "購入先目安", "価格（税込）"],
        ["M5Stack Basic V2.7", "Amazon", "¥7,990"],
        ["MAX31855 ブレークアウト", "秋月電子", "¥980"],
        ["K型熱電対 (1m)", "Amazon 等", "¥700"],
        ["MicroSD (16GB)", "量販店", "¥600"],
        ["ジャンパワイヤ", "秋月電子", "¥200"],
        ["合計", "—", "¥9,570"],
    ]
    tbl_cost = make_table(
        s,
        Inches(0.35), Inches(1.5), Inches(5.8), cost_rows,
        [Inches(2.9), Inches(1.5), Inches(1.4)],
        row_height=Inches(0.40)
    )
    # 合計行（行6）を C_ORANGE 着色
    color_row(tbl_cost, 6, C_ORANGE, C_WHITE, bold=True)

    # 右ブロック見出し
    add_textbox(s, "市販温度ロガーとの価格比較",
                Inches(6.5), Inches(1.1), Inches(6.48), Inches(0.35),
                size=Pt(15), color=C_NAVY, bold=True, font=FONT_JP)

    # 市販品比較テーブル
    market_rows = [
        ["カテゴリ", "代表製品", "概算価格"],
        ["本ツール（自作）", "M5Stack + MAX31855", "¥9,570"],
        ["標準製品", "USB データロガー", "¥30,000〜¥50,000"],
        ["プロフェッショナル", "産業用データロガー", "¥100,000〜¥250,000+"],
    ]
    tbl_mkt = make_table(
        s,
        Inches(6.5), Inches(1.5), Inches(6.48), market_rows,
        [Inches(2.8), Inches(2.18), Inches(1.5)],
        row_height=Inches(0.40)
    )
    # 本ツール行（行1）を C_NAVY 着色
    color_row(tbl_mkt, 1, C_NAVY, C_WHITE, bold=True)

    # コスト削減強調ボックス
    add_rect(s, Inches(0.35), Inches(4.35), Inches(12.6), Inches(1.4), C_NAVY)
    add_textbox(s,
                "標準製品の 1/3.1 の価格で同等機能を実現",
                Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.6),
                size=Pt(22), color=C_WHITE, bold=True, font=FONT_JP,
                align=PP_ALIGN.CENTER)
    add_textbox(s,
                "本ツール ¥9,570 vs 標準ロガー ¥30,000〜 → 最大 ¥20,000 以上削減",
                Inches(0.5), Inches(5.15), Inches(12.3), Inches(0.5),
                size=Pt(15), color=C_ORANGE, font=FONT_JP,
                align=PP_ALIGN.CENTER)
    add_textbox(s,
                "（SD 自動記録 ・ LCD リアルタイム表示 ・ 上下限アラーム機能 完全実装）",
                Inches(0.5), Inches(5.70), Inches(12.3), Inches(0.35),
                size=Pt(12), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)

    # 機能チェックリスト
    add_rect(s, Inches(0.35), Inches(5.85), Inches(12.6), Inches(1.0), C_GRAY_LT)
    add_multiline_textbox(s, [
        "✓ SD 自動記録  ✓ LCD リアルタイム表示  ✓ 上下限アラーム（ビープ音）",
        "✓ スタンドアロン動作  ✓ EEPROM 設定保存  ✓ 統計解析（平均・σ・Max/Min）",
    ],
        Inches(0.55), Inches(5.92), Inches(12.2), Inches(0.82),
        size=Pt(12), color=C_DARK_TXT, font=FONT_JP,
        align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# ▼ スライド 5 — システム全体像
# ---------------------------------------------------------------------------
def slide05(prs):
    s = _blank_slide(prs)
    header_bar(s, "システム全体像", "Section 1  プロジェクト概要")
    footer_bar(s, 5)

    # 左ブロック（入力）
    add_rect(s, Inches(0.35), Inches(1.15), Inches(3.5), Inches(5.5), C_GRAY_LT)
    add_textbox(s, "入力",
                Inches(0.45), Inches(1.25), Inches(3.3), Inches(0.4),
                size=Pt(14), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "K型熱電対",
        "↓",
        "MAX31855",
        "↓",
        "GPIO5 (CS)",
        "",
        "電源:",
        "USB-C or",
        "バッテリー",
    ],
        Inches(0.5), Inches(1.75), Inches(3.2), Inches(4.6),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP, align=PP_ALIGN.CENTER)

    # 矢印1
    add_textbox(s, "→",
                Inches(3.9), Inches(3.5), Inches(0.25), Inches(0.5),
                size=Pt(24), color=C_ORANGE, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    # 中央ブロック（M5Stack）
    add_rect(s, Inches(4.2), Inches(1.15), Inches(5.0), Inches(5.5), C_NAVY)
    add_textbox(s, "M5Stack Basic V2.7",
                Inches(4.3), Inches(1.25), Inches(4.8), Inches(0.5),
                size=Pt(16), color=C_WHITE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)
    add_multiline_textbox(s, [
        "ESP32 (Xtensa LX6 × 2, 240MHz)",
        "RAM:   520 KB SRAM",
        "Flash: 16 MB",
        "LCD:   320×240 TFT (ILI9342)",
        "EEPROM: 4 KB (仮想)",
        "Speaker 内蔵",
    ],
        Inches(4.4), Inches(1.85), Inches(4.6), Inches(4.5),
        size=Pt(13), color=C_WHITE, font=FONT_JP, align=PP_ALIGN.CENTER)

    # 矢印2
    add_textbox(s, "→",
                Inches(9.25), Inches(3.5), Inches(0.25), Inches(0.5),
                size=Pt(24), color=C_ORANGE, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    # 右ブロック（出力）
    add_rect(s, Inches(9.55), Inches(1.15), Inches(3.2), Inches(5.5), C_GRAY_LT)
    add_textbox(s, "出力",
                Inches(9.65), Inches(1.25), Inches(3.0), Inches(0.4),
                size=Pt(14), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "MicroSD",
        "→ CSV ファイル",
        "",
        "LCD",
        "→ リアルタイム表示",
        "",
        "Speaker",
        "→ アラーム音",
    ],
        Inches(9.7), Inches(1.75), Inches(2.9), Inches(4.6),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# ▼ スライド 6 — ハードウェア構成
# ---------------------------------------------------------------------------
def slide06(prs):
    s = _blank_slide(prs)
    header_bar(s, "ハードウェア構成", "Section 2  ハードウェア")
    footer_bar(s, 6)

    rows = [
        ["部品名", "型番・仕様", "役割"],
        ["メイン基板", "M5Stack Core2 (ESP32)", "メイン制御・LCD・ボタン・Speaker"],
        ["温度センサー変換", "Adafruit MAX31855", "K型熱電対→デジタル変換 (SPI)"],
        ["熱電対", "K型 (∅1.0mm 等)", "温度計測 (−200〜+1350°C)"],
        ["ストレージ", "MicroSD カード (FAT32)", "CSV データ記録"],
        ["不揮発メモリ", "ESP32 内蔵 EEPROM 仮想 (4KB)", "アラーム設定値の保存"],
        ["電源", "USB-C 5V or Lipo バッテリー", "電源供給"],
        ["接続", "ジャンパワイヤ + ブレッドボード", "センサー配線（ハンダ不要）"],
    ]
    make_table(s, CX, Inches(1.15), Inches(12.6), rows,
               [Inches(2.5), Inches(4.5), Inches(5.6)],
               row_height=Inches(0.5))
    add_textbox(s,
                "※ MicroSD ジャケットは M5Stack 本体内蔵。MAX31855 は別途 SPI 接続ブレークアウト基板を使用。",
                Inches(0.35), Inches(5.65), Inches(12.6), Inches(0.4),
                size=Pt(12), color=C_GRAY_MID, italic=True, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 7 — 配線・ピン割当
# ---------------------------------------------------------------------------
def slide07(prs):
    s = _blank_slide(prs)
    header_bar(s, "配線・ピン割当", "Section 2  ハードウェア")
    footer_bar(s, 7)

    add_textbox(s, "SPI バス共有（LCD + MAX31855）",
                Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.35),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)

    rows = [
        ["GPIO", "信号名", "接続先", "備考"],
        ["GPIO 18", "SCK (SPI Clock)", "MAX31855 CLK", "LCD と共有"],
        ["GPIO 19", "MISO (SPI Data)", "MAX31855 DO", "LCD と共有"],
        ["GPIO 23", "MOSI (SPI Data out)", "(LCD のみ)", "MAX31855 は MOSI 不要"],
        ["GPIO 5",  "CS (Chip Select)",  "MAX31855 CS", "MAX31855_CS 定数"],
        ["GPIO 4",  "CS (SD Card)",      "MicroSD CS",  "TFCARD_CS_PIN 定数"],
    ]
    make_table(s, CX, Inches(1.55), Inches(12.6), rows,
               [Inches(2.0), Inches(2.5), Inches(4.0), Inches(4.1)],
               row_height=Inches(0.45))

    # 警告ボックス
    add_rect(s, Inches(0.35), Inches(4.3), Inches(12.6), Inches(1.05),
             C_RED_LT, C_RED, line_width=1.5)
    add_textbox(s,
                "⚠ SD カードの CS は GPIO4 (TFCARD_CS_PIN) を明示指定すること",
                Inches(0.55), Inches(4.35), Inches(12.2), Inches(0.4),
                size=Pt(13), color=C_RED, bold=True, font=FONT_JP)
    add_textbox(s,
                "SD.begin(TFCARD_CS_PIN, SPI, 40000000);  // ← GPIO5 デフォルトでは MAX31855 と衝突！",
                Inches(0.55), Inches(4.78), Inches(12.2), Inches(0.45),
                size=Pt(12), color=C_DARK_TXT, font=FONT_CODE)

    add_rect(s, CX, Inches(5.45), Inches(12.6), Inches(0.7), C_GRAY_LT)
    add_textbox(s,
                "EEPROM: ESP32 内蔵フラッシュを仮想 EEPROM として使用（外部配線なし）",
                Inches(0.55), Inches(5.52), Inches(12.2), Inches(0.55),
                size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 8 — ソフトウェア概要
# ---------------------------------------------------------------------------
def slide08(prs):
    s = _blank_slide(prs)
    header_bar(s, "ソフトウェアアーキテクチャ", "Section 3  ソフトウェア設計")
    footer_bar(s, 8)

    rows = [
        ["ファイル", "クラス / 関数", "責務"],
        ["main.cpp", "setup() / loop()", "起動処理・3タスクのスケジューリング"],
        ["Tasks.cpp", "IO_Task()", "センサ読取・ボタン検出・アラーム判定 (10ms)"],
        ["Tasks.cpp", "Logic_Task()", "状態遷移・Welford統計演算 (50ms)"],
        ["Tasks.cpp", "UI_Task()", "LCD 描画（差分更新） (200ms)"],
        ["SDManager.cpp", "SDManager クラス", "CSV ファイル作成・ヘッダ/データ書込"],
        ["EEPROMManager.cpp", "EEPROMManager クラス", "アラーム値の読書き・チェックサム検証"],
        ["MeasurementCore.cpp", "MeasurementCore クラス", "純粋ロジック（ユニットテスト用）"],
    ]
    make_table(s, CX, Inches(1.15), Inches(12.6), rows,
               [Inches(2.8), Inches(3.0), Inches(6.8)],
               row_height=Inches(0.48))

    add_rect(s, CX, Inches(5.45), Inches(12.6), Inches(1.0), C_GRAY_LT)
    add_multiline_textbox(s, [
        "・「単一 loop()」ではなく millis() による周期分割で処理の干渉を防止",
        "・グローバル構造体 GlobalData (G) で各タスク間のデータを共有（PLC の D/M デバイス設計に着想）",
    ],
        Inches(0.55), Inches(5.52), Inches(12.2), Inches(0.88),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 9 — 状態機械
# ---------------------------------------------------------------------------
def slide09(prs):
    s = _blank_slide(prs)
    header_bar(s, "状態機械（4 状態）", "Section 3  ソフトウェア設計")
    footer_bar(s, 9)

    cards = [
        (C_NAVY,   "IDLE",         "待機中 / 温度監視",  "起動直後 / RESULT 後\n/ 設定終了後"),
        (C_ORANGE, "RUN",          "計測中 / CSV 記録",  "IDLE → BtnA"),
        (C_GREEN,  "RESULT",       "統計結果表示 2ページ", "RUN → BtnA"),
        (C_PURPLE, "ALARM\nSETTING", "HI/LO 閾値設定",   "IDLE → BtnB"),
    ]
    cw, ch = Inches(2.8), Inches(3.0)
    xs = [Inches(0.35), Inches(3.45), Inches(6.55), Inches(9.65)]

    for i, (fill, title, desc, cond) in enumerate(cards):
        add_rect(s, xs[i], Inches(1.15), cw, ch, fill)
        add_textbox(s, title,
                    xs[i] + Inches(0.1), Inches(1.3), cw - Inches(0.2), Inches(0.6),
                    size=Pt(18), color=C_WHITE, bold=True, font=FONT_EN,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, desc,
                    xs[i] + Inches(0.1), Inches(1.95), cw - Inches(0.2), Inches(0.55),
                    size=Pt(12), color=C_WHITE, font=FONT_JP,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, cond,
                    xs[i] + Inches(0.1), Inches(2.55), cw - Inches(0.2), Inches(0.5),
                    size=Pt(11), color=C_ORANGE if i == 0 else C_WHITE, font=FONT_JP,
                    align=PP_ALIGN.CENTER)
        if i < 3:
            add_textbox(s, "▶",
                        xs[i] + cw + Inches(0.1), Inches(2.5),
                        Inches(0.25), Inches(0.5),
                        size=Pt(16), color=C_GRAY_MID, font=FONT_EN,
                        align=PP_ALIGN.CENTER)

    rows = [
        ["遷移元", "遷移先", "トリガー", "備考"],
        ["IDLE",          "RUN",          "BtnA 短押し",      "統計値・SD ファイルを初期化"],
        ["RUN",           "RESULT",       "BtnA 短押し",      "SD ファイルをクローズ"],
        ["RESULT",        "IDLE",         "BtnA 短押し",      "—"],
        ["IDLE",          "ALARM_SETTING","BtnB 短押し",      "SettingIndex=0 (HI側) から開始"],
        ["ALARM_SETTING", "IDLE",         "BtnA 短押し × 2","LO 確定時に EEPROM 保存"],
    ]
    make_table(s, CX, Inches(4.3), Inches(12.6), rows,
               [Inches(2.0), Inches(2.0), Inches(4.0), Inches(4.6)],
               row_height=Inches(0.44))


# ---------------------------------------------------------------------------
# ▼ スライド 10 — ボタン操作一覧
# ---------------------------------------------------------------------------
def slide10(prs):
    s = _blank_slide(prs)
    header_bar(s, "ボタン操作一覧", "Section 4  操作マニュアル")
    footer_bar(s, 10)

    add_textbox(s,
                "M5Stack 前面の物理ボタン 3 個（左から BtnA / BtnB / BtnC）。"
                "短押し（エッジ検出）のみ。長押し機能なし。",
                CX, Inches(1.15), Inches(12.6), Inches(0.5),
                size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    rows = [
        ["状態",           "BtnA (左)",         "BtnB (中)",         "BtnC (右)"],
        ["IDLE",           "→ RUN（計測開始）", "→ ALARM_SETTING",  "（なし）"],
        ["RUN",            "→ RESULT（計測終了）", "（なし）",        "（なし）"],
        ["RESULT",         "→ IDLE（リセット）", "ページ切替 1/2↔2/2", "（なし）"],
        ["ALARM HI",       "→ LO 設定へ移行",  "+5°C",              "−5°C"],
        ["ALARM LO",       "→ IDLE（EEPROM 保存）", "+5°C",          "−5°C"],
    ]
    make_table(s, CX, Inches(1.75), Inches(12.6), rows,
               [Inches(2.0), Inches(3.4), Inches(3.4), Inches(3.8)],
               row_height=Inches(0.46))

    add_textbox(s,
                "※ ALARM_SETTING は IDLE からのみ進入可。RUN / RESULT 中はボタン非対応。",
                CX, Inches(5.1), Inches(12.6), Inches(0.4),
                size=Pt(12), color=C_GRAY_MID, italic=True, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 11 — IDLE 状態
# ---------------------------------------------------------------------------
def slide11(prs):
    s = _blank_slide(prs)
    header_bar(s, "操作マニュアル ① — IDLE（待機）", "Section 4  操作マニュアル")
    footer_bar(s, 11)

    # LCD モック
    lcd_mock(s, Inches(0.35), Inches(1.15), Inches(6.2), Inches(5.0), [
        "STATE: IDLE",
        "",
        "Temp:",
        "  25.3  C",
        "",
        "SD Ready",
        "",
        "[BtnA] Start   [BtnB] Setting",
    ], C_GRAY_MID)

    # 右側説明
    add_textbox(s, "IDLE 状態の動作仕様",
                Inches(6.8), Inches(1.15), Inches(6.1), Inches(0.45),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)

    rows = [
        ["項目", "内容"],
        ["温度更新周期", "500ms（TC_READ_INTERVAL_MS）"],
        ["フィルタ", "LPF  α=0.1（1次遅れ）"],
        ["異常時表示", "---.-- C（センサ未接続 / NaN）"],
        ["アラーム表示", "温度テキストを赤 (HI) / 青 (LO) に変色"],
        ["SD 状態表示", "GREEN=Ready, RED=Error, YELLOW=Not Ready"],
        ["起動リトライ", "センサ初期化を最大 5 回自動リトライ"],
    ]
    make_table(s, Inches(6.8), Inches(1.75), Inches(6.1), rows,
               [Inches(2.2), Inches(3.9)],
               row_height=Inches(0.47))


# ---------------------------------------------------------------------------
# ▼ スライド 12 — RUN 状態
# ---------------------------------------------------------------------------
def slide12(prs):
    s = _blank_slide(prs)
    header_bar(s, "操作マニュアル ② — RUN（計測中）", "Section 4  操作マニュアル")
    footer_bar(s, 12)

    lcd_mock(s, Inches(0.35), Inches(1.15), Inches(6.2), Inches(5.0), [
        "STATE: RUN",
        "",
        "Temp:",
        "  540.2  C",
        "",
        "Samples:   142",
        "SD: /DATA_0000.csv",
        "",
        "[BtnA] Stop / Reset",
    ], C_ORANGE)

    add_textbox(s, "RUN 状態の動作仕様",
                Inches(6.8), Inches(1.15), Inches(6.1), Inches(0.45),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)

    rows = [
        ["項目", "内容"],
        ["計測開始時", "統計値をリセット、CSV を DATA_0000.csv から採番"],
        ["統計更新", "Logic_Task 50ms 毎に Welford 法で逐次計算"],
        ["CSV 書込", "10 サンプル毎 (IO_Task の SD_WRITE_INTERVAL)"],
        ["アラーム", "HI/LO 超過で温度を赤/青表示 + ビープ音"],
        ["Samples 行", "G.D_Count = Logic_Task の積算回数"],
        ["SD ファイル", "/DATA_0000.csv（起動後の連番 0 始まり）"],
    ]
    make_table(s, Inches(6.8), Inches(1.75), Inches(6.1), rows,
               [Inches(2.5), Inches(3.6)],
               row_height=Inches(0.47))


# ---------------------------------------------------------------------------
# ▼ スライド 13 — RESULT 状態
# ---------------------------------------------------------------------------
def slide13(prs):
    s = _blank_slide(prs)
    header_bar(s, "操作マニュアル ③ — RESULT（結果表示）",
               "Section 4  操作マニュアル")
    footer_bar(s, 13)

    # ページ 1/2
    lcd_mock(s, Inches(0.35), Inches(1.15), Inches(5.8), Inches(5.0), [
        "STATE: RESULT (1/2)",
        "",
        "Temp:",
        "  540.2  C",
        "",
        "Avg:",
        "  539.8  C",
        "",
        "[BtnA] Reset  [BtnB] Next",
    ], C_GREEN)

    add_textbox(s, "BtnB ▶",
                Inches(6.2), Inches(3.3), Inches(0.9), Inches(0.5),
                size=Pt(16), color=C_ORANGE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    # ページ 2/2
    lcd_mock(s, Inches(7.2), Inches(1.15), Inches(5.8), Inches(5.0), [
        "STATE: RESULT (2/2)",
        "",
        "StdDev:",
        "   0.8  C",
        "Range:   10.0  C",
        "Max:    550.0  C",
        "Min:    530.0  C",
        "",
        "[BtnA] Reset  [BtnB] Prev",
    ], C_GREEN)

    add_textbox(s,
                "BtnB でページを切り替え | SD ファイルは RUN→RESULT 遷移時に自動クローズ済み",
                CX, Inches(6.25), Inches(12.6), Inches(0.5),
                size=Pt(12), color=C_GRAY_MID, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 14 — ALARM_SETTING 状態
# ---------------------------------------------------------------------------
def slide14(prs):
    s = _blank_slide(prs)
    header_bar(s, "操作マニュアル ④ — ALARM_SETTING（アラーム設定）",
               "Section 4  操作マニュアル")
    footer_bar(s, 14)

    lcd_mock(s, Inches(0.35), Inches(1.15), Inches(5.8), Inches(5.0), [
        "STATE: HI_ALARM SETTING",
        "",
        "Current:",
        "  600.0  C",
        "",
        "[BtnA] Next",
        "[BtnB] +5C  [BtnC] -5C",
    ], C_PURPLE)

    add_textbox(s, "BtnA ▶",
                Inches(6.2), Inches(3.3), Inches(0.9), Inches(0.5),
                size=Pt(16), color=C_ORANGE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    lcd_mock(s, Inches(7.2), Inches(1.15), Inches(5.8), Inches(5.0), [
        "STATE: LO_ALARM SETTING",
        "",
        "Current:",
        "  400.0  C",
        "",
        "[BtnA] Save & Exit",
        "[BtnB] +5C  [BtnC] -5C",
    ], C_PURPLE)

    add_textbox(s,
                "① IDLE で BtnB → ② HI 設定 → ③ BtnA で LO へ → "
                "④ LO 設定 → ⑤ BtnA で EEPROM 保存・IDLE 復帰",
                CX, Inches(6.3), Inches(12.6), Inches(0.5),
                size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 15 — CSV ファイル仕様
# ---------------------------------------------------------------------------
def slide15(prs):
    s = _blank_slide(prs)
    header_bar(s, "データ管理 — CSV ファイル仕様", "Section 5  データ管理")
    footer_bar(s, 15)

    add_textbox(s,
                "ファイル名: /DATA_0000.csv（0 始まり連番）　"
                "形式: UTF-8 テキスト、CRLF 改行、先頭ヘッダ行あり",
                CX, Inches(1.15), Inches(12.6), Inches(0.4),
                size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    rows = [
        ["#", "列名（ヘッダ）", "型・単位", "内容"],
        ["1",  "ElapsedSec", "uint32 / 秒",    "RUN 開始からの経過秒数 (millis/1000)"],
        ["2",  "Temp_C",     "float / °C",     "LPF フィルタ後の温度値"],
        ["3",  "State",      "文字列",          "現在の状態名（\"RUN\" 固定）"],
        ["4",  "Samples",    "uint32",         "累積サンプル数（Welford のカウント）"],
        ["5",  "Average_C",  "float / °C",     "累積平均温度"],
        ["6",  "StdDev_C",   "float / °C",     "累積標準偏差"],
        ["7",  "Max_C",      "float / °C",     "計測開始からの最大温度"],
        ["8",  "Min_C",      "float / °C",     "計測開始からの最小温度"],
        ["9",  "HI_ALARM",   "\"true\"/\"false\"", "上限アラームフラグ"],
        ["10", "LO_ALARM",   "\"true\"/\"false\"", "下限アラームフラグ"],
    ]
    make_table(s, CX, Inches(1.7), Inches(12.6), rows,
               [Inches(0.8), Inches(2.3), Inches(2.0), Inches(7.5)],
               row_height=Inches(0.44))

    add_textbox(s,
                "例: 60,540.2,RUN,120,539.8,0.8,550.0,530.0,false,false",
                CX, Inches(6.35), Inches(12.0), Inches(0.4),
                size=Pt(11), color=C_GRAY_MID, font=FONT_CODE)


# ---------------------------------------------------------------------------
# ▼ スライド 16 — Excel 活用方法
# ---------------------------------------------------------------------------
def slide16(prs):
    s = _blank_slide(prs)
    header_bar(s, "データ管理 — Excel での活用", "Section 5  データ管理")
    footer_bar(s, 16)

    rows = [
        ["手順", "操作", "備考"],
        ["①", "BtnA を押して RESULT 状態へ",            "SD ファイルが自動クローズされる"],
        ["②", "MicroSD を PC へ挿入",                   "カードリーダー経由"],
        ["③", "DATA_0000.csv を Excel で開く",          "[ファイル→開く] でカンマ区切り選択"],
        ["④", "ElapsedSec 列を横軸に折れ線グラフ",      "Temp_C, Average_C を Y 軸に追加"],
        ["⑤", "HI_ALARM = \"true\" の行を条件付き書式", "赤背景でハイライト表示"],
    ]
    make_table(s, CX, Inches(1.15), Inches(12.6), rows,
               [Inches(0.8), Inches(4.5), Inches(7.3)],
               row_height=Inches(0.5))

    add_rect(s, CX, Inches(4.8), Inches(12.6), Inches(1.6), C_GRAY_LT)
    add_multiline_textbox(s, [
        "・StdDev と Range を見ることで温度ばらつきの傾向が分かる",
        "・ファイルは毎計測ごとに連番で新規作成されるため上書きによるデータ消失なし",
        "・NaN 値（センサ異常時）は Excel フィルタで除外可能",
    ],
        Inches(0.55), Inches(4.9), Inches(12.2), Inches(1.4),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 17 — タスク設計
# ---------------------------------------------------------------------------
def slide17(prs):
    s = _blank_slide(prs)
    header_bar(s, "技術詳細 — タスク設計（3 タスク周期分割）",
               "Section 6  技術詳細")
    footer_bar(s, 17)

    add_textbox(s,
                "loop() 内で millis() のオーバーフロー安全な計算（unsigned 演算）で "
                "3 つの独立した周期タスクを実現。スレッドなし・RTOS なし。",
                CX, Inches(1.15), Inches(12.6), Inches(0.45),
                size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    tasks = [
        ("IO_Task",    "10 ms",  "95%", [
            "MAX31855 センサ読取",
            "（実際は 500ms 毎）",
            "LPF フィルタ適用 (α=0.1)",
            "BtnA/B/C エッジ検出",
            "ヒステリシス付きアラーム判定",
            "SD カード書込 (10サンプル毎)",
        ]),
        ("Logic_Task", "50 ms",  "98%", [
            "Welford 法による逐次統計更新",
            "BtnA/B/C イベント処理",
            "状態遷移制御",
            "EEPROM 保存処理",
        ]),
        ("UI_Task",    "200 ms", "97.5%", [
            "状態別 LCD 描画",
            "差分更新（全画面消去は遷移時のみ）",
            "RESULT 2 ページ描画",
            "ALARM_SETTING HI/LO 描画",
        ]),
    ]

    cw = Inches(3.9)
    xs = [Inches(0.35), Inches(4.55), Inches(8.75)]

    for i, (title, period, margin, items) in enumerate(tasks):
        cx = xs[i]
        # カード枠
        add_rect(s, cx, Inches(1.7), cw, Inches(4.6),
                 C_WHITE, C_NAVY, line_width=1.5)
        # オレンジ上端アクセントバー
        add_rect(s, cx, Inches(1.7), cw, Inches(0.07), C_ORANGE)
        # タイトル
        add_textbox(s, title, cx + Inches(0.1), Inches(1.82),
                    cw - Inches(0.2), Inches(0.5),
                    size=Pt(18), color=C_NAVY, bold=True, font=FONT_EN,
                    align=PP_ALIGN.CENTER)
        # 周期
        add_textbox(s, period, cx + Inches(0.1), Inches(2.35),
                    cw - Inches(0.2), Inches(0.4),
                    size=Pt(14), color=C_ORANGE, bold=True, font=FONT_EN,
                    align=PP_ALIGN.CENTER)
        # 余裕率
        add_textbox(s, f"実行余裕: {margin}", cx + Inches(0.1), Inches(2.75),
                    cw - Inches(0.2), Inches(0.3),
                    size=Pt(12), color=C_GREEN, font=FONT_JP,
                    align=PP_ALIGN.CENTER)
        # 処理内容
        add_multiline_textbox(s, ["・" + t for t in items],
                              cx + Inches(0.15), Inches(3.12),
                              cw - Inches(0.3), Inches(2.9),
                              size=Pt(12), color=C_DARK_TXT, font=FONT_JP)

    add_textbox(s,
                "IO_Task は 10ms で起動するが、センサ実読取は TC_READ_INTERVAL_MS = 500ms ごとに制限"
                "（MAX31855 変換完了待ち）",
                CX, Inches(6.4), Inches(12.6), Inches(0.4),
                size=Pt(12), color=C_NAVY, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 18 — センサ読取とフィルタ
# ---------------------------------------------------------------------------
def slide18(prs):
    s = _blank_slide(prs)
    header_bar(s, "技術詳細 — センサ読取と LPF フィルタ",
               "Section 6  技術詳細")
    footer_bar(s, 18)

    # 左ブロック
    add_textbox(s, "センサ読取シーケンス",
                Inches(0.35), Inches(1.15), Inches(5.9), Inches(0.45),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)
    rows = [
        ["項目", "値/内容"],
        ["IC", "Adafruit MAX31855"],
        ["通信", "ハードウェア SPI (SPI.h)"],
        ["変換時間", "最大 100ms (データシート)"],
        ["読取間隔", "500ms (TC_READ_INTERVAL_MS)"],
        ["失敗時", "最大 3 回リトライ (readThermocouple())"],
    ]
    make_table(s, Inches(0.35), Inches(1.65), Inches(5.9), rows,
               [Inches(2.2), Inches(3.7)],
               row_height=Inches(0.47))

    # 右ブロック
    add_textbox(s, "LPF（1次遅れフィルタ）",
                Inches(6.55), Inches(1.15), Inches(6.1), Inches(0.45),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)

    add_rect(s, Inches(6.55), Inches(1.65), Inches(6.1), Inches(1.2), C_GRAY_LT)
    add_multiline_textbox(s, [
        "y[n] = y[n-1] × (1 − α) + x[n] × α",
        "α = FILTER_ALPHA = 0.1",
    ],
        Inches(6.7), Inches(1.75), Inches(5.8), Inches(1.0),
        size=Pt(14), color=C_DARK_TXT, font=FONT_CODE,
        align=PP_ALIGN.CENTER)

    rows2 = [
        ["項目", "値"],
        ["収束サンプル数 (新値 90%)", "約 22 サンプル（≈ 11 秒）"],
        ["初回起動時", "raw 値をそのまま初期値に設定 (NaN 収束待ち回避)"],
        ["適用タイミング", "新データ到着時のみ（同値での重複適用なし）"],
    ]
    make_table(s, Inches(6.55), Inches(2.95), Inches(6.1), rows2,
               [Inches(2.5), Inches(3.6)],
               row_height=Inches(0.47))

    add_textbox(s,
                "α=0.1 は実機検証と CSV データ分析で決定。値が大きいほど応答速、小さいほど平滑",
                CX, Inches(6.3), Inches(12.6), Inches(0.4),
                size=Pt(12), color=C_GRAY_MID, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 19 — Welford 統計アルゴリズム
# ---------------------------------------------------------------------------
def slide19(prs):
    s = _blank_slide(prs)
    header_bar(s, "技術詳細 — Welford 法による逐次統計",
               "Section 6  技術詳細")
    footer_bar(s, 19)

    # 左
    add_textbox(s, "Welford 法とは",
                Inches(0.35), Inches(1.15), Inches(7.6), Inches(0.4),
                size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)
    add_rect(s, Inches(0.35), Inches(1.6), Inches(7.6), Inches(1.8), C_GRAY_LT)
    add_textbox(s,
                "全データを保存せずに「1 サンプルずつ受け取るたびに」"
                "平均・分散・標準偏差を更新するアルゴリズム。"
                "メモリ O(1) で長時間計測に最適。",
                Inches(0.5), Inches(1.7), Inches(7.3), Inches(1.6),
                size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    add_rect(s, Inches(0.35), Inches(3.5), Inches(7.6), Inches(2.2),
             C_WHITE, C_NAVY, line_width=1.0)
    add_textbox(s, "更新式（Logic_Task 内）",
                Inches(0.5), Inches(3.58), Inches(7.3), Inches(0.38),
                size=Pt(13), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "1. delta = x[n] - prevMean",
        "2. Sum  += x[n]",
        "3. newMean = Sum / Count",
        "4. delta2 = x[n] - newMean",
        "5. M2 += delta * delta2",
        "6. StdDev = sqrt(M2 / Count)",
    ],
        Inches(0.5), Inches(4.0), Inches(7.2), Inches(1.6),
        size=Pt(12), color=C_DARK_TXT, font=FONT_CODE)

    # 右
    add_rect(s, Inches(8.3), Inches(1.15), Inches(4.7), Inches(3.0),
             C_WHITE, C_GREEN, line_width=1.5)
    add_textbox(s, "採用理由",
                Inches(8.45), Inches(1.25), Inches(4.4), Inches(0.4),
                size=Pt(15), color=C_GREEN, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・メモリ使用量 O(1)（計測時間に依存しない）",
        "・数値安定性: double 型で丸め誤差を最小化",
        "・リアルタイム更新可能（50ms 毎に再計算）",
        "・RAM 使用率: 7.2% を計測時間によらず維持",
    ],
        Inches(8.45), Inches(1.75), Inches(4.4), Inches(2.2),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    add_rect(s, Inches(8.3), Inches(4.25), Inches(4.7), Inches(1.8), C_GRAY_LT)
    add_textbox(s, "実装ポイント",
                Inches(8.45), Inches(4.33), Inches(4.4), Inches(0.38),
                size=Pt(13), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・D_M2, D_Sum は double 型（精度確保）",
        "・最終値は handleButtonA() で RUN→RESULT",
        "  遷移時にも再計算",
    ],
        Inches(8.45), Inches(4.72), Inches(4.4), Inches(1.2),
        size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 20 — アラーム判定アルゴリズム
# ---------------------------------------------------------------------------
def slide20(prs):
    s = _blank_slide(prs)
    header_bar(s, "技術詳細 — ヒステリシス付きアラーム判定",
               "Section 6  技術詳細")
    footer_bar(s, 20)

    add_rect(s, CX, Inches(1.15), Inches(12.6), Inches(0.55), C_GRAY_LT)
    add_textbox(s,
                "デフォルト: HI = 600.0°C ／ LO = 400.0°C ／ ヒステリシス = 5.0°C"
                "  ※ EEPROM 保存・電源 OFF 後も保持",
                Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.4),
                size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    rows = [
        ["種別", "トリガー条件", "クリア条件", "音", "LCD 色"],
        ["HI アラーム", "温度 ≥ HI しきい値", "温度 < HI − 5°C", "2kHz, 500ms", "赤 (RED)"],
        ["LO アラーム", "温度 ≤ LO しきい値", "温度 > LO + 5°C", "1kHz, 500ms", "青 (BLUE)"],
        ["非アラーム",  "—",                   "—",              "なし",        "白 (WHITE)"],
    ]
    make_table(s, CX, Inches(1.8), Inches(12.6), rows,
               [Inches(1.5), Inches(3.5), Inches(3.5), Inches(1.7), Inches(2.4)],
               row_height=Inches(0.5))

    add_rect(s, CX, Inches(4.0), Inches(6.8), Inches(2.5), C_GRAY_LT)
    add_textbox(s, "ヒステリシスの効果",
                Inches(0.5), Inches(4.1), Inches(6.5), Inches(0.4),
                size=Pt(14), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "温度   601°C → [HI 発報] ビープ",
        "       600°C → [継続]（まだクリアされない）",
        "       596°C → [継続]（< 595°C まで待機）",
        "       594°C → [HI クリア]",
    ],
        Inches(0.5), Inches(4.55), Inches(6.5), Inches(1.75),
        size=Pt(13), color=C_DARK_TXT, font=FONT_CODE)

    add_rect(s, Inches(7.5), Inches(4.0), Inches(5.5), Inches(2.5),
             C_WHITE, C_NAVY, line_width=1.0)
    add_textbox(s, "実装のポイント",
                Inches(7.65), Inches(4.1), Inches(5.2), Inches(0.4),
                size=Pt(14), color=C_NAVY, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・チラつき防止: しきい値近辺での誤発報を抑制",
        "・NaN 値ガード: isnan() で無効値をスキップ",
        "・updateAlarmFlags() は独立関数",
        "  → ユニットテスト可能",
    ],
        Inches(7.65), Inches(4.55), Inches(5.2), Inches(1.75),
        size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 21 — デバッグ対応事例
# ---------------------------------------------------------------------------
def slide21(prs):
    s = _blank_slide(prs)
    header_bar(s, "品質・テスト — デバッグ対応事例",
               "Section 7  品質・テスト")
    footer_bar(s, 21)

    rows = [
        ["#", "症状", "原因", "対策"],
        ["01", "温度が全て 0.0°C になる",
         "SD.begin() が GPIO5 既定 → MAX31855 と CS 衝突",
         "SD.begin(TFCARD_CS_PIN, SPI, 40000000) で GPIO4 を明示"],
        ["02", "起動直後にアラームが誤発報",
         "setup() 中の IO_Task 呼び出しが閾値を超える値を一時読取",
         "setup() 末尾でアラームフラグを強制リセット"],
        ["03", "RUN 中の CSV に Average=0.0 が続く",
         "D_Count<2 のタイミングで書込 → Welford 未完了",
         "SD 書込条件を D_Count≥10 に変更"],
        ["04", "SD 書込後に温度が NaN になる",
         "SPI バスが SD 書込後に解放されない",
         "writeData() 後に SPI.end(); SPI.begin() を実行"],
        ["05", "フィルタが NaN にはまったまま収束しない",
         "初回読取 NaN を前回値に適用してしまう",
         "初回は raw 値をそのまま代入（isnan チェック追加）"],
    ]
    make_table(s, CX, Inches(1.15), Inches(12.6), rows,
               [Inches(0.5), Inches(3.5), Inches(3.5), Inches(5.1)],
               row_height=Inches(0.5))

    add_rect(s, CX, Inches(5.85), Inches(12.6), Inches(0.7), C_GRAY_LT)
    add_textbox(s,
                "共通パターン: ハードウェアリソース（SPI/GPIO）の競合は CS 指定の明示化で解決。"
                "ソフトウェア状態は単体テスト可能な関数に分離することで早期発見。",
                Inches(0.55), Inches(5.93), Inches(12.2), Inches(0.58),
                size=Pt(13), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ スライド 22 — テスト実績
# ---------------------------------------------------------------------------
def slide22(prs):
    s = _blank_slide(prs)
    header_bar(s, "品質・テスト — テスト実績", "Section 7  品質・テスト")
    footer_bar(s, 22)

    add_textbox(s,
                "PlatformIO native テストにより 33 件すべて合格（test_measurement_core.cpp）",
                CX, Inches(1.15), Inches(12.6), Inches(0.45),
                size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    rows = [
        ["テスト対象", "テストケース概要", "件数", "結果"],
        ["状態遷移",     "IDLE→RUN→RESULT→IDLE の連続遷移が仕様通り動作",  "5 件",  "✓ 全合格"],
        ["統計計算",     "既知値での Average / StdDev が期待値に一致",      "8 件",  "✓ 全合格"],
        ["アラーム判定", "ヒステリシス境界値での HI/LO フラグ挙動",         "6 件",  "✓ 全合格"],
        ["EEPROM 操作",  "書込→読込の一致、NaN/Inf の拒否",               "4 件",  "✓ 全合格"],
        ["フィルタ",     "初回 NaN 時の raw 代入, α=0.1 の収束特性",       "4 件",  "✓ 全合格"],
        ["CSV 出力",     "NaN データ行の \"NaN\" 文字列出力",              "3 件",  "✓ 全合格"],
        ["ボタン操作",   "ALARM_SETTING の BtnB +5°C / BtnC -5°C",       "3 件",  "✓ 全合格"],
    ]
    tbl = make_table(s, CX, Inches(1.8), Inches(9.0), rows,
                     [Inches(1.8), Inches(4.5), Inches(0.75), Inches(1.95)],
                     row_height=Inches(0.46))

    # 合格率カード
    add_rect(s, Inches(9.5), Inches(5.45), Inches(3.13), Inches(1.15), C_GREEN)
    add_textbox(s, "33 / 33",
                Inches(9.5), Inches(5.52), Inches(3.13), Inches(0.65),
                size=Pt(28), color=C_WHITE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)
    add_textbox(s, "100%合格",
                Inches(9.5), Inches(6.17), Inches(3.13), Inches(0.35),
                size=Pt(14), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# ▼ スライド 23 — まとめ
# ---------------------------------------------------------------------------
def slide23(prs):
    s = _blank_slide(prs)
    header_bar(s, "まとめ", "Summary")
    footer_bar(s, 23)

    # 成果ボックス
    add_rect(s, CX, Inches(1.15), Inches(12.6), Inches(2.4), C_NAVY)
    add_textbox(s, "実現したこと",
                Inches(0.55), Inches(1.22), Inches(12.2), Inches(0.4),
                size=Pt(18), color=C_WHITE, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・PLC 設計手法（D/M デバイス）を C++ に応用した状態機械の実装",
        "・3 タスク（10 / 50 / 200 ms）のノン・プリエンプティブ周期スケジューリング",
        "・Welford 法によるリアルタイム統計（RAM O(1)、計測時間に無依存）",
        "・SPI バス競合・CS 衝突などハードウェア固有バグの特定と対策",
        "・native テスト 33 件 100% 合格で品質担保",
    ],
        Inches(0.55), Inches(1.68), Inches(12.2), Inches(1.75),
        size=Pt(14), color=C_WHITE, font=FONT_JP)

    # 学習ポイントボックス
    add_rect(s, CX, Inches(3.7), Inches(5.9), Inches(2.8),
             C_WHITE, C_ORANGE, line_width=1.5)
    add_textbox(s, "学んだこと",
                Inches(0.55), Inches(3.8), Inches(5.5), Inches(0.4),
                size=Pt(16), color=C_ORANGE, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・ハードウェアとソフトウェアの境界を意識した設計の重要性",
        "・デバッグは再現性・ログ出力から始める",
        "・「小さなテスト可能な関数」に分割することで不具合が減る",
        "・NaN / 初期化の扱いは想定外の動作の温床になる",
    ],
        Inches(0.55), Inches(4.28), Inches(5.7), Inches(2.1),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    # 今後の展望ボックス
    add_rect(s, Inches(6.6), Inches(3.7), Inches(6.05), Inches(2.8),
             C_WHITE, C_GREEN, line_width=1.5)
    add_textbox(s, "次のステップ",
                Inches(6.75), Inches(3.8), Inches(5.7), Inches(0.4),
                size=Pt(16), color=C_GREEN, bold=True, font=FONT_JP)
    add_multiline_textbox(s, [
        "・RTC（DS3231）追加でファイル名を YYYYMMDD_HHMMSS 形式に",
        "・多チャネル化（同時 8ch 計測）",
        "・WiFi / MQTT でクラウド送信",
        "・FreeRTOS 移行でより精密なタスク管理",
    ],
        Inches(6.75), Inches(4.28), Inches(5.7), Inches(2.1),
        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)


# ---------------------------------------------------------------------------
# ▼ メイン
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    funcs = [
        slide01, slide02, slide03, slide04, slide05,
        slide06, slide07, slide08, slide09, slide10,
        slide11, slide12, slide13, slide14, slide15,
        slide16, slide17, slide18, slide19, slide20,
        slide21, slide22, slide23,
    ]

    for i, fn in enumerate(funcs):
        print(f"  Building Slide {i+1:02d} — {fn.__name__} ...")
        fn(prs)

    assert len(prs.slides) == TOTAL_SLIDES, (
        f"スライド枚数が一致しません: {len(prs.slides)} != {TOTAL_SLIDES}"
    )

    out_path = os.path.join(os.path.dirname(__file__),
                            "Simple_Temp_Tool_Overview_v2.pptx")
    prs.save(out_path)
    print(f"\n✅ 保存完了: {out_path}")
    print(f"   スライド枚数: {len(prs.slides)} 枚")


if __name__ == "__main__":
    main()
