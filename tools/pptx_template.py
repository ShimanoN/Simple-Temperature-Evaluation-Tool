#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
pptx_template.py
----------------
Universal PowerPoint generation template.
Copy this file to docs/<project>/ and rename it to create_presentation.py.

Workflow:
  1. Edit ProjectConfig (single source of truth for all variable values)
  2. Add slide functions following the pattern below
  3. Run:  python create_presentation.py

Design principles:
  - ProjectConfig centralizes ALL variable values (hardware names, costs, etc.)
    This prevents the "Core2 vs Basic V2.7" class of drift bugs.
  - calc_table_bottom() computes the exact y-coordinate where a table ends,
    preventing overlap bugs when placing blocks below tables.
  - English comments throughout (AI-readable); Japanese strings in output.
  - Coordinate system: origin (0,0) at top-left, y increases downward, unit = Inches.

Slide layout (16:9, 13.33" × 7.5"):
  ┌──────────────────────────────────────── 13.33" ──┐
  │  Header bar  (y=0.0 .. 1.0")                     │ ← HDR_H = 1.0"
  ├──────────────────────────────────────────────────┤
  │  Content area (y=1.1 .. 6.85")                   │ ← CY=1.1", CH=5.75"
  │    x starts at CX=0.35"                          │
  │                                                  │
  ├──────────────────────────────────────────────────┤
  │  Footer bar  (y=6.9 .. 7.5")                     │ ← FTR_Y=6.9"
  └──────────────────────────────────────────────────┘
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ===========================================================================
# ▼ PROJECT CONFIG  ← Edit this section for each new project
#    Single source of truth: every variable value lives here ONLY.
#    Never hardcode project-specific strings inside slide functions.
# ===========================================================================

class ProjectConfig:
    # ---- Basic info --------------------------------------------------------
    TITLE           = "プロジェクトタイトルをここに"
    SUBTITLE        = "サブタイトル（ハードウェア名 × センサー名 等）"
    TECH_STACK      = "ESP32  /  C++ (Arduino)  /  PlatformIO"
    VERSION         = "v1.0.0  ｜  2026 年 X 月"

    # ---- Hardware ----------------------------------------------------------
    # IMPORTANT: Define hardware names here once.
    # If the hardware changes, update ONLY this block.
    HARDWARE_MCU    = "M5Stack Basic V2.7"      # e.g. "M5Stack Basic V2.7"
    HARDWARE_SENSOR = "MAX31855"
    HARDWARE_MCU_PRICE   = "¥7,990"
    HARDWARE_MCU_SOURCE  = "Amazon"

    # ---- Cost breakdown rows (for cost slide) ------------------------------
    # Format: [component_name, source, price_string]
    COST_ROWS = [
        ["部品", "購入先目安", "価格（税込）"],
        ["M5Stack Basic V2.7",     "Amazon",   "¥7,990"],
        ["MAX31855 ブレークアウト", "秋月電子", "¥980"],
        ["K型熱電対 (1m)",         "Amazon 等","¥700"],
        ["MicroSD (16GB)",         "量販店",   "¥600"],
        ["ジャンパワイヤ",         "秋月電子", "¥200"],
        ["合計",                   "—",        "¥9,570"],
    ]
    COST_TOTAL      = "¥9,570"

    # ---- Output path -------------------------------------------------------
    OUTPUT_DIR  = os.path.join(os.path.dirname(__file__), ".")
    OUTPUT_FILE = "output_presentation.pptx"


# ===========================================================================
# ▼ COLOR CONSTANTS  (customize per brand/project)
# ===========================================================================
C_NAVY     = RGBColor(0x0D, 0x47, 0x7F)
C_ORANGE   = RGBColor(0xFF, 0x6D, 0x00)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_GRAY_LT  = RGBColor(0xF4, 0xF4, 0xF4)
C_GRAY_MID = RGBColor(0x99, 0x99, 0x99)
C_RED_LT   = RGBColor(0xFF, 0xE0, 0xE0)
C_BLACK    = RGBColor(0x00, 0x00, 0x00)

# ===========================================================================
# ▼ FONT CONSTANTS
# ===========================================================================
FONT_JP   = "游ゴシック"
FONT_EN   = "Segoe UI"
FONT_CODE = "Consolas"

# ===========================================================================
# ▼ LAYOUT CONSTANTS  (16:9 slide, 13.33" × 7.5")
# ===========================================================================
W      = Inches(13.33)   # slide width
H      = Inches(7.5)     # slide height
HDR_H  = Inches(1.0)     # header bar height
FTR_Y  = Inches(6.9)     # footer bar top
FTR_H  = Inches(0.6)     # footer bar height
CX     = Inches(0.35)    # content area left edge
CY     = Inches(1.1)     # content area top edge
CW     = Inches(12.63)   # content area width
CH     = Inches(5.75)    # content area height
CONTENT_BOTTOM = Inches(6.85)  # last safe y before footer

RECT = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE


# ===========================================================================
# ▼ KEY UTILITY: COORDINATE CALCULATION
#    Use these to avoid overlap bugs.
# ===========================================================================

def calc_table_bottom(start_y_inches: float, num_rows: int,
                      row_height_inches: float) -> float:
    """
    Calculate the y-coordinate (in inches) where a table ends.

    Usage example:
        table_end = calc_table_bottom(1.5, 7, 0.40)   # → 4.30
        add_rect(s, Inches(0.35), Inches(table_end + 0.1), ...)  # safe margin

    Args:
        start_y_inches:    y position of the table top (in inches)
        num_rows:          total number of rows including header
        row_height_inches: height of each row (in inches)

    Returns:
        float: y position in inches where the table ends
    """
    return start_y_inches + num_rows * row_height_inches


def safe_y_after_table(start_y_inches: float, num_rows: int,
                       row_height_inches: float, margin_inches: float = 0.1) -> float:
    """
    Return a safe y-coordinate to place an element below a table.
    Adds a small margin to prevent overlap.

    Args:
        margin_inches: gap between table bottom and next element (default 0.1")
    """
    return calc_table_bottom(start_y_inches, num_rows, row_height_inches) + margin_inches


# ===========================================================================
# ▼ SLIDE HELPER FUNCTIONS
# ===========================================================================

def _blank_slide(prs: Presentation):
    """Add and return a blank (white background) slide."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color: RGBColor,
             line_color=None, line_width=None):
    """Add a filled rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1.0)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                size=Pt(13), color=None, bold=False, italic=False,
                font=FONT_JP, align=PP_ALIGN.LEFT, wrap=True):
    """Add a single-paragraph textbox."""
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_multiline_textbox(slide, lines: list, x, y, w, h,
                          size=Pt(13), color=None, bold=False,
                          font=FONT_JP, align=PP_ALIGN.LEFT,
                          line_spacing_pt: float = None):
    """
    Add a textbox with multiple paragraphs (one per line).
    Useful for bulleted or checklist-style content.
    """
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True

    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            p.line_spacing = _Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return tb


def header_bar(slide, title: str, section: str = ""):
    """
    Render the standard header bar (navy background + title + section label).
    Standard positioning: x=0, y=0, w=13.33", h=1.0"
    """
    add_rect(slide, Inches(0), Inches(0), W, HDR_H, C_NAVY)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.08), Inches(10.5), Inches(0.55),
                size=Pt(24), color=C_WHITE, bold=True, font=FONT_JP)
    if section:
        add_textbox(slide, section,
                    Inches(0.35), Inches(0.63), Inches(10.5), Inches(0.3),
                    size=Pt(12), color=C_ORANGE, font=FONT_JP)


def footer_bar(slide, page_num: int, total: int = None):
    """
    Render the standard footer bar.
    Shows project title on the left, page number on the right.
    """
    add_rect(slide, Inches(0), FTR_Y, W, FTR_H, C_NAVY)
    add_textbox(slide, ProjectConfig.TITLE,
                Inches(0.35), Inches(6.95), Inches(9.0), Inches(0.4),
                size=Pt(11), color=C_GRAY_MID, font=FONT_EN)
    page_text = f"{page_num}" if total is None else f"{page_num} / {total}"
    add_textbox(slide, page_text,
                Inches(11.5), Inches(6.95), Inches(1.5), Inches(0.4),
                size=Pt(11), color=C_WHITE, font=FONT_EN, align=PP_ALIGN.RIGHT)


def make_table(slide, x, y, w, rows: list, col_widths: list,
               row_height=Inches(0.40), header_fill=C_NAVY,
               header_text_color=C_WHITE, alt_fill=C_GRAY_LT):
    """
    Create a formatted table with a colored header row.

    IMPORTANT: To avoid overlap with elements below this table, use:
        end_y = safe_y_after_table(start_y_inches, len(rows), row_height_inches)

    Args:
        x, y, w:          position and width (Inches objects)
        rows:             list of lists; rows[0] is the header row
        col_widths:       list of Inches objects, one per column
        row_height:       Inches object for all rows
        header_fill:      RGBColor for header row background
        header_text_color: RGBColor for header row text
        alt_fill:         RGBColor for alternating data rows

    Returns:
        pptx Table object
    """
    from pptx.util import Inches as _Inches
    num_rows = len(rows)
    num_cols = len(rows[0])
    total_h = row_height * num_rows

    tbl = slide.shapes.add_table(num_rows, num_cols, x, y, w, total_h).table

    # Set column widths
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    # Set row heights and fill cells
    for ri, row_data in enumerate(rows):
        tbl.rows[ri].height = row_height
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            # Style
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                tf = cell.text_frame
                tf.paragraphs[0].runs[0].font.color.rgb = header_text_color
                tf.paragraphs[0].runs[0].font.bold = True
                tf.paragraphs[0].runs[0].font.name = FONT_JP
                tf.paragraphs[0].runs[0].font.size = Pt(12)
            else:
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_WHITE
                tf = cell.text_frame
                tf.paragraphs[0].runs[0].font.name = FONT_JP
                tf.paragraphs[0].runs[0].font.size = Pt(12)
    return tbl


def color_row(tbl, row_index: int, fill: RGBColor,
              text_color: RGBColor, bold: bool = False):
    """Apply background and text color to an entire table row."""
    row = tbl.rows[row_index]
    for cell in row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        tf = cell.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.color.rgb = text_color
                if bold:
                    run.font.bold = True


# ===========================================================================
# ▼ SLIDE FUNCTIONS  — one function per slide
#
# Naming convention: slide01(), slide02(), ... slide23()
#
# Coordinate checklist before adding any element below a table:
#   1. table_start_y  = (the y value passed to make_table, in inches)
#   2. table_num_rows = len(rows)
#   3. row_h          = (the row_height value, in inches)
#   4. next_y         = safe_y_after_table(table_start_y, table_num_rows, row_h)
#   5. Pass Inches(next_y) as the y argument for the next element
# ===========================================================================

# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def slide01(prs):
    s = _blank_slide(prs)
    _set_bg(s, C_NAVY)

    add_textbox(s, ProjectConfig.TITLE,
                Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.9),
                size=Pt(36), color=C_WHITE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    # SUBTITLE references ProjectConfig — change hardware name in ONE place only
    add_textbox(s, f"{ProjectConfig.HARDWARE_MCU} × {ProjectConfig.HARDWARE_SENSOR}  熱電対温度計測システム",
                Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.55),
                size=Pt(18), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)

    add_rect(s, Inches(1.0), Inches(3.62), Inches(11.3), Inches(0.04), C_ORANGE)

    add_textbox(s, "リアルタイム計測 ／ 統計解析 ／ SD カード記録 ／ アラーム機能",
                Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.45),
                size=Pt(15), color=C_ORANGE, font=FONT_JP,
                align=PP_ALIGN.CENTER)

    add_textbox(s, ProjectConfig.TECH_STACK,
                Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.45),
                size=Pt(14), color=C_GRAY_MID, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    add_textbox(s, ProjectConfig.VERSION,
                Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.4),
                size=Pt(13), color=C_WHITE, font=FONT_JP,
                align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — (Copy this block as a template for new slides)
# ---------------------------------------------------------------------------
def slide02(prs):
    s = _blank_slide(prs)
    header_bar(s, "スライドタイトル", "Section X  セクション名")
    footer_bar(s, 2)

    # --- Example: text block ---
    add_textbox(s, "見出しテキスト",
                CX, CY, CW, Inches(0.4),
                size=Pt(18), color=C_NAVY, bold=True)

    # --- Example: table with safe positioning below ---
    TABLE_START_Y   = 1.6    # inches — table top
    TABLE_ROW_H     = 0.40   # inches — row height
    rows = [
        ["列1", "列2", "列3"],
        ["データ A", "詳細 A", "値 A"],
        ["データ B", "詳細 B", "値 B"],
    ]
    col_widths = [Inches(4.0), Inches(4.0), Inches(4.63)]

    make_table(s, CX, Inches(TABLE_START_Y), CW, rows, col_widths,
               row_height=Inches(TABLE_ROW_H))

    # --- Place next element BELOW the table (no overlap) ---
    next_y = safe_y_after_table(TABLE_START_Y, len(rows), TABLE_ROW_H, margin_inches=0.15)
    # next_y is now guaranteed to be below the table

    add_rect(s, CX, Inches(next_y), CW, Inches(0.6), C_NAVY)
    add_textbox(s, "テーブルの下に安全に配置されたブロック",
                Inches(0.5), Inches(next_y + 0.1), Inches(12.3), Inches(0.4),
                size=Pt(14), color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)


# ===========================================================================
# ▼ MAIN — assemble all slides and save
# ===========================================================================
def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Register slide functions in order
    slide_funcs = [
        slide01,
        slide02,
        # slide03, slide04, ...  ← add more here
    ]

    total = len(slide_funcs)
    for i, func in enumerate(slide_funcs, 1):
        func(prs)
        print(f"  [{i:02d}/{total:02d}] {func.__name__} ... done")

    output_path = os.path.join(
        ProjectConfig.OUTPUT_DIR,
        ProjectConfig.OUTPUT_FILE
    )
    prs.save(output_path)
    print(f"\n✅ 保存完了: {output_path}  スライド枚数: {total} 枚")


if __name__ == "__main__":
    main()
