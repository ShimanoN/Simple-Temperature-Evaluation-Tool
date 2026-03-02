#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
spec_to_pptx.py
---------------
Domain-agnostic Markdown spec → PowerPoint converter.
Works for any presentation type: business, technical, research, recruitment, etc.

Usage:
    python tools/spec_to_pptx.py path/to/presentation_spec.md

Spec format: see tools/PPTX_CREATION_WORKFLOW.md for full documentation.

Design:
    - parse_spec()           : Parses YAML front matter + slide sections from .md
    - _interpolate_params()  : Replaces {{key}} placeholders with front matter values
    - render_*()             : One renderer per layout type (7 layouts)
    - safe_y_after_table()   : Prevents overlap bugs by calculating table bottom y
    - main()                 : Orchestrates parse → interpolate → render → save

Front matter reserved keys:
    title    : Main title (required)
    subtitle : Subtitle line on title slide
    tagline  : Orange accent line on title slide
    meta     : Gray info line (author, org, tech stack, etc.)
    date     : Date / version shown on title slide
    output   : Output .pptx file path (required)
   (any other key is accessible via {{key}} in slide content)
"""

import os
import re
import sys
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ===========================================================================
# ▼ COLOR CONSTANTS
# ===========================================================================
C_NAVY     = RGBColor(0x0D, 0x47, 0x7F)
C_ORANGE   = RGBColor(0xFF, 0x6D, 0x00)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GREEN    = RGBColor(0x27, 0xAE, 0x60)
C_RED      = RGBColor(0xC0, 0x39, 0x2B)
C_GRAY_LT  = RGBColor(0xF4, 0xF4, 0xF4)
C_GRAY_MID = RGBColor(0x99, 0x99, 0x99)
C_RED_LT   = RGBColor(0xFF, 0xE0, 0xE0)
C_BLACK    = RGBColor(0x00, 0x00, 0x00)

# ===========================================================================
# ▼ FONT CONSTANTS
# ===========================================================================
FONT_JP   = "游ゴシック"
FONT_EN   = "Segoe UI"
FONT_CODE = "Consolas"

# ===========================================================================
# ▼ LAYOUT CONSTANTS  (16:9, 13.33" × 7.5")
# ===========================================================================
W              = Inches(13.33)
H              = Inches(7.5)
HDR_H          = Inches(1.0)
FTR_Y          = Inches(6.9)
FTR_H          = Inches(0.6)
CX             = Inches(0.35)
CY             = Inches(1.1)
CW             = Inches(12.63)
CH             = Inches(5.75)
CONTENT_BOTTOM = 6.85   # inches — last safe y before footer

RECT = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE


# ===========================================================================
# ▼ COORDINATE HELPERS  — use these, never hardcode y after a table
# ===========================================================================

def calc_table_bottom(start_y: float, num_rows: int, row_h: float) -> float:
    """Return the y-coordinate (inches) where the table ends."""
    return start_y + num_rows * row_h


def safe_y_after_table(start_y: float, num_rows: int, row_h: float,
                        margin: float = 0.12) -> float:
    """Return a safe y (inches) for the first element below a table."""
    return calc_table_bottom(start_y, num_rows, row_h) + margin


def clamp_to_content(y: float, element_h: float) -> float:
    """Warn if element would overflow footer area."""
    if y + element_h > CONTENT_BOTTOM:
        print(f"  ⚠ WARNING: element at y={y:.2f}\" h={element_h:.2f}\" "
              f"exceeds content bottom ({CONTENT_BOTTOM}\"). May overlap footer.")
    return y


# ===========================================================================
# ▼ PPTX PRIMITIVE HELPERS
# ===========================================================================

def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color,
             line_color=None, line_width=None):
    shape = slide.shapes.add_shape(RECT, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1.0)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                size=Pt(13), color=None, bold=False, italic=False,
                font=FONT_JP, align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return tb


def add_multiline_textbox(slide, lines: list, x, y, w, h,
                           size=Pt(13), color=None, bold=False,
                           font=FONT_JP, align=PP_ALIGN.LEFT):
    if color is None:
        color = C_DARK_TXT
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = size
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = font
    return tb


def header_bar(slide, title: str, section: str = ""):
    add_rect(slide, Inches(0), Inches(0), W, HDR_H, C_NAVY)
    add_textbox(slide, title,
                Inches(0.35), Inches(0.08), Inches(10.5), Inches(0.55),
                size=Pt(24), color=C_WHITE, bold=True, font=FONT_JP)
    if section:
        add_textbox(slide, section,
                    Inches(0.35), Inches(0.63), Inches(10.5), Inches(0.3),
                    size=Pt(12), color=C_ORANGE, font=FONT_JP)


def footer_bar(slide, page_num: int, project_title: str = ""):
    add_rect(slide, Inches(0), FTR_Y, W, FTR_H, C_NAVY)
    add_textbox(slide, project_title,
                Inches(0.35), Inches(6.95), Inches(9.0), Inches(0.4),
                size=Pt(11), color=C_GRAY_MID, font=FONT_EN)
    add_textbox(slide, str(page_num),
                Inches(11.5), Inches(6.95), Inches(1.5), Inches(0.4),
                size=Pt(11), color=C_WHITE, font=FONT_EN, align=PP_ALIGN.RIGHT)


def _ptitle(config: dict) -> str:
    """Return the footer display title: 'title' field takes priority over 'project'."""
    return config.get('title') or config.get('project', '')


def make_table(slide, x, y, w, rows: list, col_widths: list,
               row_height=Inches(0.40), header_fill=C_NAVY,
               header_text_color=C_WHITE, alt_fill=C_GRAY_LT):
    num_rows = len(rows)
    num_cols = len(rows[0])
    total_h  = row_height * num_rows
    tbl = slide.shapes.add_table(num_rows, num_cols, x, y, w, total_h).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row_data in enumerate(rows):
        tbl.rows[ri].height = row_height
        for ci, cell_text in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                r = cell.text_frame.paragraphs[0].runs[0]
                r.font.color.rgb = header_text_color
                r.font.bold = True
                r.font.name = FONT_JP
                r.font.size = Pt(12)
            else:
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_fill
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_WHITE
                r = cell.text_frame.paragraphs[0].runs[0]
                r.font.name = FONT_JP
                r.font.size = Pt(12)
    return tbl


def color_row(tbl, row_index, fill, text_color, bold=False):
    row = tbl.rows[row_index]
    for cell in row.cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = text_color
                if bold:
                    run.font.bold = True


# ===========================================================================
# ▼ TEMPLATE VARIABLE INTERPOLATION
# ===========================================================================

def _interpolate(value: str, config: dict) -> str:
    """
    Replace {{key}} placeholders in a string with values from config.
    Unknown keys are left as-is (e.g. {{unknown}} stays {{unknown}}).
    """
    def _replace(m):
        key = m.group(1).strip()
        return str(config.get(key, f'{{{{{key}}}}}'))
    return re.sub(r'\{\{(\w+)\}\}', _replace, value)


def _interpolate_params(params: dict, config: dict) -> dict:
    """
    Walk params dict recursively and apply {{key}} substitution.
    Handles string values, lists of strings, and lists of lists (table rows).
    """
    result = {}
    for k, v in params.items():
        if isinstance(v, str):
            result[k] = _interpolate(v, config)
        elif isinstance(v, list):
            processed = []
            for item in v:
                if isinstance(item, str):
                    processed.append(_interpolate(item, config))
                elif isinstance(item, list):
                    processed.append([
                        _interpolate(cell, config) if isinstance(cell, str) else cell
                        for cell in item
                    ])
                else:
                    processed.append(item)
            result[k] = processed
        else:
            result[k] = v
    return result


# ===========================================================================
# ▼ SPEC PARSER
# ===========================================================================

def parse_spec(filepath: str):
    """
    Parse a presentation_spec.md file.

    Returns:
        config (dict):  YAML front matter (project-wide settings)
        slides (list):  list of {'layout': str, 'params': dict}
    """
    with open(filepath, 'r', encoding='utf-8-sig') as f:  # utf-8-sig strips BOM on Windows
        content = f.read()

    # --- Extract YAML front matter ---
    fm_match = re.match(r'^---\r?\n(.*?)\r?\n---\r?\n', content, re.DOTALL)
    if fm_match:
        config = yaml.safe_load(fm_match.group(1)) or {}
        body   = content[fm_match.end():]
    else:
        print("WARNING: No YAML front matter found (--- ... ---)")
        config = {}
        body   = content

    # --- Extract slide sections ---
    # Pattern: ## slide_NN | layout: LAYOUT_NAME
    # Captures everything until the next ## slide_ or end of string
    slide_header = re.compile(
        r'^## slide_\d+\s*\|\s*layout:\s*(\S+)\s*\n(.*?)(?=^## slide_|\Z)',
        re.MULTILINE | re.DOTALL
    )
    slides = []
    for match in slide_header.finditer(body):
        layout     = match.group(1).strip()
        yaml_block = match.group(2).strip()
        try:
            params = yaml.safe_load(yaml_block) if yaml_block else {}
        except yaml.YAMLError as e:
            print(f"WARNING: YAML parse error in slide (layout={layout}): {e}")
            params = {}
        slides.append({'layout': layout, 'params': params or {}})

    return config, slides


# ===========================================================================
# ▼ LAYOUT RENDERERS
# ===========================================================================

def render_title(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: title
    Builds the title slide from front matter fields.
    All fields can also be overridden per-slide in params.

    Front matter fields used:
        title    → large main title                  (required)
        subtitle → white subtitle below title        (optional)
        tagline  → orange accent line                (optional)
        meta     → gray info line (author/org/stack) (optional)
        date     → white date/version at bottom      (optional)
    """
    _set_bg(s, C_NAVY)

    title    = params.get('title')    or config.get('title',    'Presentation Title')
    subtitle = params.get('subtitle') or config.get('subtitle', '')
    tagline  = params.get('tagline')  or config.get('tagline',  '')
    meta     = params.get('meta')     or config.get('meta',     '')
    date     = params.get('date')     or config.get('date',     '')

    add_textbox(s, title,
                Inches(1.0), Inches(1.8), Inches(11.3), Inches(0.9),
                size=Pt(36), color=C_WHITE, bold=True, font=FONT_JP,
                align=PP_ALIGN.CENTER)

    if subtitle:
        add_textbox(s, subtitle,
                    Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.55),
                    size=Pt(18), color=C_WHITE, font=FONT_JP,
                    align=PP_ALIGN.CENTER)

    add_rect(s, Inches(1.0), Inches(3.62), Inches(11.3), Inches(0.04), C_ORANGE)

    if tagline:
        add_textbox(s, tagline,
                    Inches(1.0), Inches(3.8), Inches(11.3), Inches(0.45),
                    size=Pt(15), color=C_ORANGE, font=FONT_JP,
                    align=PP_ALIGN.CENTER)

    if meta:
        add_textbox(s, meta,
                    Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.45),
                    size=Pt(14), color=C_GRAY_MID, font=FONT_JP,
                    align=PP_ALIGN.CENTER)

    if date:
        add_textbox(s, date,
                    Inches(1.0), Inches(5.2), Inches(11.3), Inches(0.4),
                    size=Pt(13), color=C_WHITE, font=FONT_JP,
                    align=PP_ALIGN.CENTER)


def render_toc(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: toc
    params: title (optional), items (list of strings)
    """
    title = params.get('title', '目次')
    items = params.get('items', [])

    header_bar(s, title, "")
    footer_bar(s, page_num, _ptitle(config))
    for i, item in enumerate(items):
        y = 1.35 + i * 0.62
        # Left accent bar
        add_rect(s, Inches(0.5), Inches(y), Inches(0.06), Inches(0.42), C_ORANGE)
        add_textbox(s, item,
                    Inches(0.72), Inches(y), Inches(11.5), Inches(0.42),
                    size=Pt(18), color=C_DARK_TXT, font=FONT_JP)


def render_section_divider(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: section_divider
    params: section_num, section_title, description (optional)
    """
    _set_bg(s, C_NAVY)

    section_num   = params.get('section_num', '')
    section_title = params.get('section_title', '')
    description   = params.get('description', '')

    add_textbox(s, section_num,
                Inches(1.0), Inches(2.2), Inches(11.3), Inches(0.6),
                size=Pt(22), color=C_ORANGE, bold=True, font=FONT_EN,
                align=PP_ALIGN.CENTER)

    add_textbox(s, section_title,
                Inches(1.0), Inches(2.9), Inches(11.3), Inches(0.8),
                size=Pt(36), color=C_WHITE, bold=True, font=FONT_JP,
                align=PP_ALIGN.CENTER)

    add_rect(s, Inches(3.5), Inches(3.85), Inches(6.3), Inches(0.04), C_ORANGE)

    if description:
        add_textbox(s, description,
                    Inches(1.0), Inches(4.05), Inches(11.3), Inches(0.5),
                    size=Pt(15), color=C_GRAY_MID, font=FONT_JP,
                    align=PP_ALIGN.CENTER)

    footer_bar(s, page_num, _ptitle(config))


def render_bullet(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: bullet
    params: title, section, bullets (list), highlight (optional), sub_text (optional)
    """
    title     = params.get('title', '')
    section   = params.get('section', '')
    bullets   = params.get('bullets', [])
    highlight = params.get('highlight', '')
    sub_text  = params.get('sub_text', '')

    header_bar(s, title, section)
    footer_bar(s, page_num, _ptitle(config))

    # Bullet items
    bullet_start_y = 1.2
    for i, bullet in enumerate(bullets):
        y = bullet_start_y + i * 0.52
        add_rect(s, Inches(0.35), Inches(y + 0.12), Inches(0.06), Inches(0.25), C_ORANGE)
        add_textbox(s, bullet,
                    Inches(0.55), Inches(y), Inches(12.2), Inches(0.45),
                    size=Pt(14), color=C_DARK_TXT, font=FONT_JP)

    # Optional highlight box below bullets
    if highlight or sub_text:
        box_y = bullet_start_y + len(bullets) * 0.52 + 0.2
        box_y = clamp_to_content(box_y, 0.8)
        add_rect(s, CX, Inches(box_y), CW, Inches(0.55), C_NAVY)
        add_textbox(s, highlight,
                    Inches(0.5), Inches(box_y + 0.07), Inches(12.3), Inches(0.38),
                    size=Pt(16), color=C_WHITE, bold=True, font=FONT_JP,
                    align=PP_ALIGN.CENTER)
        if sub_text:
            sub_y = box_y + 0.62
            sub_y = clamp_to_content(sub_y, 0.4)
            add_textbox(s, sub_text,
                        Inches(0.5), Inches(sub_y), Inches(12.3), Inches(0.4),
                        size=Pt(13), color=C_DARK_TXT, font=FONT_JP,
                        align=PP_ALIGN.CENTER)


def render_table(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: table
    params: title, section, columns, col_widths (list of floats),
            rows (list of lists), row_height (float, default 0.40),
            warning (str, optional), warning_code (str, optional),
            note (str, optional)
    """
    title        = params.get('title', '')
    section      = params.get('section', '')
    columns      = params.get('columns', [])
    col_widths_f = params.get('col_widths', [])
    data_rows    = params.get('rows', [])
    row_h        = float(params.get('row_height', 0.40))
    warning      = params.get('warning', '')
    warning_code = params.get('warning_code', '')
    note         = params.get('note', '')

    header_bar(s, title, section)
    footer_bar(s, page_num, _ptitle(config))

    # Subtitle if present
    subtitle = params.get('subtitle', '')
    if subtitle:
        add_textbox(s, subtitle,
                    CX, Inches(1.1), CW, Inches(0.35),
                    size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)

    # Build full rows (header + data)
    all_rows = [columns] + [list(r) for r in data_rows]
    if col_widths_f:
        col_widths_in = [Inches(w) for w in col_widths_f]
    else:
        # Auto-equal split across full content width (12.63")
        num_cols = len(all_rows[0]) if all_rows else 1
        unit = 12.63 / num_cols
        col_widths_in = [Inches(unit)] * num_cols

    TABLE_START_Y = 1.55
    make_table(s, CX, Inches(TABLE_START_Y), CW, all_rows, col_widths_in,
               row_height=Inches(row_h))

    # Calculate safe y below table
    next_y = safe_y_after_table(TABLE_START_Y, len(all_rows), row_h, margin=0.12)

    # Optional warning box
    if warning:
        warn_h = 1.0 if warning_code else 0.52
        clamp_to_content(next_y, warn_h)
        add_rect(s, Inches(0.35), Inches(next_y), CW, Inches(warn_h),
                 C_RED_LT, C_RED, line_width=1.5)
        add_textbox(s, warning,
                    Inches(0.55), Inches(next_y + 0.06), Inches(12.2), Inches(0.38),
                    size=Pt(13), color=C_RED, bold=True, font=FONT_JP)
        if warning_code:
            add_textbox(s, warning_code,
                        Inches(0.55), Inches(next_y + 0.48), Inches(12.2), Inches(0.4),
                        size=Pt(12), color=C_DARK_TXT, font=FONT_CODE)
        next_y = next_y + warn_h + 0.12

    # Optional note box
    if note:
        note_h = 0.52
        clamp_to_content(next_y, note_h)
        add_rect(s, CX, Inches(next_y), CW, Inches(note_h), C_GRAY_LT)
        add_textbox(s, note,
                    Inches(0.55), Inches(next_y + 0.07), Inches(12.2), Inches(0.38),
                    size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


def render_dual_table(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: dual_table
    Side-by-side two-table layout with optional highlight box and checklist.
    Fully configurable labels and highlight rows — domain-agnostic.

    params:
        title              : slide title (header bar)
        section            : section label (header bar, optional)
        left_label         : label above left table
        left_rows          : list of rows (first row = header)
        left_col_widths    : list of floats in inches (default: auto-split 5.8")
        left_highlight_last: true → highlight last row (e.g. total row) in orange
        right_label        : label above right table
        right_rows         : list of rows (first row = header)
        right_col_widths   : list of floats in inches (default: auto-split 6.48")
        right_highlight_row: int → highlight this row index in navy (0-indexed: 0=header, 1=first data row, 2=second data row, …)
        highlight          : text for navy summary box (optional)
        sub_highlight      : orange subtext in summary box (optional)
        checklist          : list of strings for gray box below (optional)

    Layout (y-axis, auto-calculated):
      y=1.10  Left/right table labels
      y=1.50  Two tables
      y=safe  Navy highlight box → orange sub_highlight → gray checklist
    """
    title               = params.get('title', '')
    section             = params.get('section', '')
    left_label          = params.get('left_label', '')
    left_rows           = params.get('left_rows', [])
    left_col_widths_f   = params.get('left_col_widths', [])
    left_highlight_last = params.get('left_highlight_last', False)
    right_label         = params.get('right_label', '')
    right_rows          = params.get('right_rows', [])
    right_col_widths_f  = params.get('right_col_widths', [])
    right_highlight_row = params.get('right_highlight_row', None)
    highlight           = params.get('highlight', '')
    sub_highlight       = params.get('sub_highlight', '')
    checklist           = params.get('checklist', [])

    header_bar(s, title, section)
    footer_bar(s, page_num, _ptitle(config))

    LEFT_W      = 5.8
    RIGHT_W     = 6.48
    LEFT_X      = 0.35
    RIGHT_X     = 6.5
    ROW_H       = 0.40
    TABLE_START = 1.50

    # --- Left table ---
    if left_label:
        add_textbox(s, left_label,
                    Inches(LEFT_X), Inches(1.1), Inches(LEFT_W), Inches(0.35),
                    size=Pt(15), color=C_NAVY, bold=True, font=FONT_JP)
    if left_rows:
        num_cols = len(left_rows[0])
        if left_col_widths_f:
            lcw = [Inches(w) for w in left_col_widths_f]
        else:
            unit = LEFT_W / num_cols
            lcw = [Inches(unit)] * num_cols
        tbl_left = make_table(s, Inches(LEFT_X), Inches(TABLE_START), Inches(LEFT_W),
                              left_rows, lcw, row_height=Inches(ROW_H))
        if left_highlight_last and len(left_rows) > 1:
            color_row(tbl_left, len(left_rows) - 1, C_ORANGE, C_WHITE, bold=True)

    # --- Right table ---
    if right_label:
        add_textbox(s, right_label,
                    Inches(RIGHT_X), Inches(1.1), Inches(RIGHT_W), Inches(0.35),
                    size=Pt(15), color=C_NAVY, bold=True, font=FONT_JP)
    if right_rows:
        num_cols = len(right_rows[0])
        if right_col_widths_f:
            rcw = [Inches(w) for w in right_col_widths_f]
        else:
            unit = RIGHT_W / num_cols
            rcw = [Inches(unit)] * num_cols
        tbl_right = make_table(s, Inches(RIGHT_X), Inches(TABLE_START), Inches(RIGHT_W),
                               right_rows, rcw, row_height=Inches(ROW_H))
        if right_highlight_row is not None and len(right_rows) > right_highlight_row:
            color_row(tbl_right, right_highlight_row, C_NAVY, C_WHITE, bold=True)

    # --- Safe y below the taller table ---
    left_end  = calc_table_bottom(TABLE_START, len(left_rows),  ROW_H)
    right_end = calc_table_bottom(TABLE_START, len(right_rows), ROW_H)
    box_y = max(left_end, right_end) + 0.12

    # --- Navy highlight box ---
    if highlight:
        hl_h = 1.0
        box_y = clamp_to_content(box_y, hl_h)
        add_rect(s, Inches(0.35), Inches(box_y), Inches(12.6), Inches(hl_h), C_NAVY)
        add_textbox(s, highlight,
                    Inches(0.5), Inches(box_y + 0.07), Inches(12.3), Inches(0.45),
                    size=Pt(20), color=C_WHITE, bold=True, font=FONT_JP,
                    align=PP_ALIGN.CENTER)
        if sub_highlight:
            add_textbox(s, sub_highlight,
                        Inches(0.5), Inches(box_y + 0.56), Inches(12.3), Inches(0.38),
                        size=Pt(14), color=C_ORANGE, font=FONT_JP,
                        align=PP_ALIGN.CENTER)
        box_y = box_y + hl_h + 0.1

    # --- Gray checklist box ---
    if checklist:
        list_h = 0.2 + len(checklist) * 0.35
        box_y = clamp_to_content(box_y, list_h)
        add_rect(s, Inches(0.35), Inches(box_y), Inches(12.6), Inches(list_h), C_GRAY_LT)
        add_multiline_textbox(s, checklist,
                              Inches(0.55), Inches(box_y + 0.07),
                              Inches(12.2), Inches(list_h - 0.1),
                              size=Pt(12), color=C_DARK_TXT, font=FONT_JP,
                              align=PP_ALIGN.CENTER)


def render_two_column(prs, s, config: dict, params: dict, page_num: int):
    """
    layout: two_column
    params: title, section, left_title, left_items (list),
            right_title, right_items (list), note (optional)
    """
    title       = params.get('title', '')
    section     = params.get('section', '')
    left_title  = params.get('left_title', '')
    left_items  = params.get('left_items', [])
    right_title = params.get('right_title', '')
    right_items = params.get('right_items', [])
    note        = params.get('note', '')

    header_bar(s, title, section)
    footer_bar(s, page_num, _ptitle(config))

    # Divider line down the center
    add_rect(s, Inches(6.62), Inches(1.15), Inches(0.04), Inches(5.55), C_GRAY_MID)

    def render_column(x_start, col_title, items):
        add_textbox(s, col_title,
                    Inches(x_start), Inches(1.2), Inches(5.9), Inches(0.38),
                    size=Pt(16), color=C_NAVY, bold=True, font=FONT_JP)
        for i, item in enumerate(items):
            y = 1.7 + i * 0.52
            add_rect(s, Inches(x_start), Inches(y + 0.12),
                     Inches(0.06), Inches(0.25), C_ORANGE)
            add_textbox(s, item,
                        Inches(x_start + 0.18), Inches(y), Inches(5.7), Inches(0.45),
                        size=Pt(13), color=C_DARK_TXT, font=FONT_JP)

    render_column(0.35, left_title, left_items)
    render_column(6.72, right_title, right_items)

    if note:
        # Place note below the longer column's last item
        max_items = max(len(left_items), len(right_items))
        note_y = 1.7 + max_items * 0.52 + 0.2
        note_y = clamp_to_content(note_y, 0.48)
        add_rect(s, CX, Inches(note_y), CW, Inches(0.48), C_GRAY_LT)
        add_textbox(s, note,
                    Inches(0.55), Inches(note_y + 0.07), Inches(12.2), Inches(0.35),
                    size=Pt(12), color=C_DARK_TXT, font=FONT_JP)


# ===========================================================================
# ▼ LAYOUT REGISTRY — maps layout name to renderer function
# ===========================================================================

LAYOUT_RENDERERS = {
    'title':            render_title,
    'toc':              render_toc,
    'section_divider':  render_section_divider,
    'bullet':           render_bullet,
    'table':            render_table,
    'dual_table':       render_dual_table,
    'two_column':       render_two_column,
}


# ===========================================================================
# ▼ MAIN
# ===========================================================================

def main():
    if len(sys.argv) < 2:
        print("Usage: python tools/spec_to_pptx.py path/to/presentation_spec.md")
        sys.exit(1)

    spec_path = sys.argv[1]
    if not os.path.exists(spec_path):
        print(f"ERROR: spec file not found: {spec_path}")
        sys.exit(1)

    print(f"📖 Reading spec: {spec_path}")
    config, slides = parse_spec(spec_path)

    if not slides:
        print("ERROR: No slides found in spec. Check ## slide_NN | layout: headers.")
        sys.exit(1)

    _display_title = config.get('title') or config.get('project', '(unnamed)')
    print(f"   Title      : {_display_title}")
    print(f"   Slide count: {len(slides)}")
    print()

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, slide_def in enumerate(slides, 1):
        layout = slide_def['layout']
        # Apply {{key}} template variable substitution before rendering
        params = _interpolate_params(slide_def['params'], config)
        renderer = LAYOUT_RENDERERS.get(layout)

        if renderer is None:
            print(f"  [{i:02d}] ⚠ Unknown layout '{layout}' — skipping")
            continue

        s = _blank_slide(prs)
        renderer(prs, s, config, params, i)
        print(f"  [{i:02d}/{len(slides):02d}] layout={layout:<20} ... done")

    # --- Determine output path ---
    output_file = config.get('output', 'output_presentation.pptx')
    spec_dir    = os.path.dirname(os.path.abspath(spec_path))
    output_path = os.path.join(spec_dir, output_file)

    prs.save(output_path)
    print(f"\n✅ 保存完了: {output_path}")
    print(f"   スライド枚数: {len(slides)} 枚")


if __name__ == "__main__":
    main()
